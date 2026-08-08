import streamlit as st
from google import genai
from google.genai import types
import gspread
from google.oauth2.service_account import Credentials
from collections import defaultdict
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import io, json, re
from datetime import date, datetime

today_str = date.today().strftime("%Y%m%d")

st.title("📋 レジメン情報抽出・登録")
st.caption("PDFをアップロードしてAIが自動解析→スプレッドシートに登録→パワポ生成まで一気通貫")
st.divider()

# ===== 認証・初期化 =====
@st.cache_resource
def get_gemini_client():
    api_key = st.secrets["gemini"]["api_key"]
    return genai.Client(api_key=api_key)

@st.cache_resource
def get_gspread_client():
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]
    creds = Credentials.from_service_account_info(
        st.secrets["gcp_service_account"], scopes=scopes
    )
    return gspread.authorize(creds)

@st.cache_resource
def get_spreadsheet():
    gc = get_gspread_client()
    return gc.open_by_url(st.secrets["spreadsheet"]["url"])

@st.cache_data(ttl=600)
def load_definition():
    gc = get_gspread_client()
    SPREADSHEET_ID = "1dLEUYSZlrIK1uHqEtEAfS1jSAPpXCIiAiAk_iaRuY-8"
    sh = gc.open_by_key(SPREADSHEET_ID)
    ws = sh.worksheet("抽出定義書")
    return ws.acell("A1").value

@st.cache_data(ttl=300)
def load_master_data():
    sh = get_spreadsheet()
    master_data = sh.worksheet("薬品マスタ").get_all_records()
    notes_data  = sh.worksheet("注意事項").get_all_records()
    return master_data, notes_data

# ===== ヘルパー関数 =====
def get_val(d, *keys, default=""):
    for key in keys:
        if key in d and d[key] is not None:
            return d[key]
    return default

def get_drugs(data):
    return data.get("drug_info") or data.get("drugs") or []

def get_basic(data):
    return data.get("basic_info") or {}

def to_half_kana(text):
    table = {
        'ア':'ｱ','イ':'ｲ','ウ':'ｳ','エ':'ｴ','オ':'ｵ',
        'カ':'ｶ','キ':'ｷ','ク':'ｸ','ケ':'ｹ','コ':'ｺ',
        'サ':'ｻ','シ':'ｼ','ス':'ｽ','セ':'ｾ','ソ':'ｿ',
        'タ':'ﾀ','チ':'ﾁ','ツ':'ﾂ','テ':'ﾃ','ト':'ﾄ',
        'ナ':'ﾅ','ニ':'ﾆ','ヌ':'ﾇ','ネ':'ﾈ','ノ':'ﾉ',
        'ハ':'ﾊ','ヒ':'ﾋ','フ':'ﾌ','ヘ':'ﾍ','ホ':'ﾎ',
        'マ':'ﾏ','ミ':'ﾐ','ム':'ﾑ','メ':'ﾒ','モ':'ﾓ',
        'ヤ':'ﾔ','ユ':'ﾕ','ヨ':'ﾖ',
        'ラ':'ﾗ','リ':'ﾘ','ル':'ﾙ','レ':'ﾚ','ロ':'ﾛ',
        'ワ':'ﾜ','ヲ':'ｦ','ン':'ﾝ',
        'ァ':'ｧ','ィ':'ｨ','ゥ':'ｩ','ェ':'ｪ','ォ':'ｫ',
        'ッ':'ｯ','ャ':'ｬ','ュ':'ｭ','ョ':'ｮ',
        'ガ':'ｶﾞ','ギ':'ｷﾞ','グ':'ｸﾞ','ゲ':'ｹﾞ','ゴ':'ｺﾞ',
        'ザ':'ｻﾞ','ジ':'ｼﾞ','ズ':'ｽﾞ','ゼ':'ｾﾞ','ゾ':'ｿﾞ',
        'ダ':'ﾀﾞ','ヂ':'ﾁﾞ','ヅ':'ﾂﾞ','デ':'ﾃﾞ','ド':'ﾄﾞ',
        'バ':'ﾊﾞ','ビ':'ﾋﾞ','ブ':'ﾌﾞ','ベ':'ﾍﾞ','ボ':'ﾎﾞ',
        'パ':'ﾊﾟ','ピ':'ﾋﾟ','プ':'ﾌﾟ','ペ':'ﾍﾟ','ポ':'ﾎﾟ',
        'ー':'ｰ','ヴ':'ｳﾞ','・':'･',
    }
    result = ''
    for char in str(text):
        result += table.get(char, char)
    return result

def shorten_regimen_name(regimen_name):
    name = regimen_name
    name = re.sub(r'[（(][^）)]*[）)]', '', name)
    for word in [
        '術前','術後','周術期','切除不能','再発','進行',
        '維持','補助','一次治療','二次治療','初回','難治性',
        '肺癌','胃癌','大腸癌','乳癌','膵癌','肝癌',
        '食道癌','子宮癌','卵巣癌','前立腺癌','膀胱癌',
        '腎癌','甲状腺癌','悪性リンパ腫','白血病',
        '骨髄腫','中皮腫','胸腺癌','胸腺腫','神経内分泌腫瘍',
    ]:
        name = name.replace(word, '')
    name = re.sub(r'\s+', ' ', name).strip('　 ')
    return name

def parse_days_num(day_str):
    days = []
    for part in str(day_str).split('|'):
        part = part.strip()
        if '-' in part:
            try:
                s, e = part.split('-')
                days.extend(range(int(s), int(e)+1))
            except:
                pass
        elif part.isdigit():
            days.append(int(part))
    return sorted(set(days))

def format_dose_text(drug):
    try:
        try:
            _raw_dose = str(drug.get('投与量数値', 0) or 0)
            dose = float(re.sub(r'[^0-9.]', '', _raw_dose) or 0)
        except (ValueError, TypeError):
            dose = 0
    except:
        dose = 0
    unit_input = str(drug.get('投与単位', '')).strip()
    dose_base  = str(drug.get('用量根拠', ''))
    v_to_mg    = drug.get('1V当たりmg', '')
    if unit_input.upper() == 'V':
        if v_to_mg != '' and str(v_to_mg).strip() != '':
            try:
                dose = dose * float(v_to_mg)
                unit_input = 'mg'
            except:
                pass
    if unit_input in ('mg','mg/body','mg/ body'):
        display_unit = 'mg'
    elif unit_input == '':
        display_unit = '' if dose_base == 'AUC依存' else 'mg'
    else:
        display_unit = unit_input
    dose_str = str(int(dose)) if dose == int(dose) else str(dose)
    return dose_str, display_unit

INJECTION_ORDER = {
    'NK1':1,'5HT3':2,'ステロイド':3,'G-CSF':4,
    '利尿薬':5,'解毒薬':6,'抗アレルギー':7,
    'H2ブロッカー':8,'電解質補正':9,'その他注射':10
}

def make_support_line(drug):
    name = to_half_kana(
        str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
    )
    dose_str, unit_str = format_dose_text(drug)
    day = str(drug.get('投与Day文字',''))
    return f"{name} {dose_str}{unit_str}({day})"

def get_regimen(protocol_no, basic_data, drug_data, master_data):
    master_dict = {m['管理コード']: m for m in master_data}
    basic = [b for b in basic_data if b['プロトコールNo'] == protocol_no]
    if not basic:
        return None
    basic = basic[0]
    drugs_raw = sorted(
        [d for d in drug_data if d['プロトコールNo'] == protocol_no],
        key=lambda x: (
            0 if str(x['投与順序']).isdigit() else 1,
            int(x['投与順序']) if str(x['投与順序']).isdigit() else 99
        )
    )
    drugs = []
    for drug in drugs_raw:
        code   = str(drug['管理コード'])
        master = master_dict.get(code, {})
        merged = dict(drug)
        merged.update({
            '一般名（全角）'        : master.get('一般名（全角）',''),
            '一般名（半角カナ）'    : master.get('一般名（半角カナ）',''),
            '採用商品名（全角）'    : master.get('採用商品名（全角）',''),
            '採用商品名（半角カナ）': master.get('採用商品名（半角カナ）',''),
            '薬効分類'             : master.get('薬効分類',''),
            '薬剤区分'             : master.get('薬剤区分',''),
            '支持療法分類'         : master.get('支持療法分類',''),
            '薬品マスタ単位'       : master.get('単位',''),
            '投与経路'             : master.get('投与経路',''),
            '1V当たりmg'           : master.get('1V当たりmg',''),
            '患者向け説明'         : master.get('患者向け説明',''),
        })
        drugs.append(merged)
    return {'basic': basic, 'drugs': drugs, 'master_dict': master_dict}

# ===== create_pptx =====
def create_pptx(protocol_no, basic_data, drug_data,
                master_data, notes_data):
    result = get_regimen(protocol_no, basic_data, drug_data, master_data)
    if result is None:
        return None
    basic       = result['basic']
    drugs       = result['drugs']
    master_dict = result['master_dict']
    cycle       = int(basic['1コース日数'])

    COLOR_HEADER  = RGBColor(0x40,0x40,0x40)
    COLOR_ICI     = RGBColor(0x70,0xAD,0x47)
    COLOR_MOL     = RGBColor(0xFF,0xE0,0x99)
    COLOR_NAUSEA  = RGBColor(0x84,0xAC,0xD4)
    COLOR_CYTO    = RGBColor(0xF4,0xB1,0x83)
    COLOR_AUX     = RGBColor(0x9D,0xD7,0xEA)
    COLOR_REST_BG = RGBColor(0xBF,0xBF,0xBF)
    COLOR_WHITE   = RGBColor(0xFF,0xFF,0xFF)
    COLOR_BLACK   = RGBColor(0x00,0x00,0x00)
    TYPE_COLOR_MAP = {
        '免疫チェックポイント阻害薬': COLOR_ICI,
        '分子標的薬': COLOR_MOL,
        '吐き気止め': COLOR_NAUSEA,
        '細胞障害性抗がん薬': COLOR_CYTO,
        '補助薬': COLOR_AUX,
        '輸液': COLOR_AUX,
        'その他支持療法': COLOR_NAUSEA,
        '造血因子': COLOR_AUX,
    }
    A4_W = Cm(25.4); A4_H = Cm(19.05)

    def rgb_hex(rgb):
        return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

    def set_cell_bg(cell, rgb):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', rgb_hex(rgb))

    def set_cell_text_multi(cell, lines, font_size=Pt(11),
                             bold=False, color=COLOR_BLACK,
                             align=PP_ALIGN.CENTER):
        tf = cell.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = str(line); run.font.size = font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold = bold; run.font.color.rgb = color

    def set_cell_margin(cell, top=Cm(0.05), bottom=Cm(0.05),
                         left=Cm(0.1), right=Cm(0.1)):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        tcPr.set('marT',str(int(top))); tcPr.set('marB',str(int(bottom)))
        tcPr.set('marL',str(int(left))); tcPr.set('marR',str(int(right)))

    def merge_cells_vertical(table, col, start_row, end_row):
        for r in range(start_row, end_row+1):
            tc = table.cell(r,col)._tc; tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:vMerge')): tcPr.remove(old)
            vMerge = etree.SubElement(tcPr, qn('a:vMerge'))
            if r == start_row: vMerge.set('val','restart')

    def add_textbox(slide, x, y, w, h, text,
                    font_size=Pt(11), bold=False,
                    color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text=str(text); run.font.size=font_size
        run.font.name = "BIZ UDゴシック"
        run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def add_textbox_multi(slide, x, y, w, h, lines,
                          font_size=Pt(11), bold=False,
                          color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text=str(line); run.font.size=font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def safe_float(val):
        try:
            return float(val) if str(val).strip() != '' else 0
        except:
            return 0

    seal_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='外来手帳シール'
         and str(n.get('プロトコールNo',''))=='共通'],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )
    tegaki_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='手帳シール但し書き'
         and str(n.get('プロトコールNo','')) in ('共通',protocol_no)],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )

    schedule_drugs = [d for d in drugs if str(d.get('④説明書',''))=='○']
    inj_schedule   = [d for d in schedule_drugs if str(d.get('投与順序',''))!='内服']
    oral_schedule  = [d for d in schedule_drugs if str(d.get('投与順序',''))=='内服']

    rp_groups = defaultdict(list)
    for drug in inj_schedule:
        rp_groups[str(drug.get('投与順序',''))].append(drug)
    sorted_rps = sorted(rp_groups.keys(),
                        key=lambda x: int(x) if x.isdigit() else 99)

    all_invest_days = set()
    for rp in sorted_rps:
        all_invest_days.update(
            parse_days_num(rp_groups[rp][0].get('投与Day数値',''))
        )
    invest_days = sorted(all_invest_days)
    columns = []; prev = 0
    for d in invest_days:
        if d > prev+1:
            rs=prev+1; re=d-1
            columns.append({'type':'休薬中',
                             'label':f'{rs}日目' if rs==re else f'{rs}〜{re}日目',
                             'days':list(range(rs,re+1))})
        columns.append({'type':'投与','label':f'{d}日目','days':[d]})
        prev = d
    if prev < cycle:
        rs=prev+1; re=cycle
        columns.append({'type':'休薬','label':f'〜{re}日目',
                        'days':list(range(rs,re+1))})

    need_bw  = any(str(d.get('用量根拠','')) in ('BW依存','AUC依存')
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')
    need_bsa = any(str(d.get('用量根拠',''))=='BSA依存'
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')

    rp_all_groups = defaultdict(list)
    for drug in drugs:
        rp = str(drug.get('投与順序',''))
        if rp != '内服':
            rp_all_groups[rp].append(drug)

    total_min_raw = sum(
        max((safe_float(d.get('投与時間数値', 0))
             for d in rp_drugs), default=0)
        for rp_drugs in rp_all_groups.values()
    ) * 60
    import math
    total_min = int(math.ceil(total_min_raw / 5) * 5)

    def get_rp_info(rp_drugs):
        rp_sorted = sorted(rp_drugs,
                           key=lambda d: INJECTION_ORDER.get(str(d.get('支持療法分類','')),99))
        names = [
            str(d.get('採用商品名（全角）','') or d.get('商品名',''))
            for d in rp_sorted
            if not str(d.get('管理コード','')).strip().startswith('IV')
        ]
        name_text = '\n'.join([n for n in names if n])
        dose_parts = []
        for d in rp_sorted:
            dose_base = str(d.get('用量根拠',''))
            kubun     = str(d.get('薬剤区分',''))
            try:
                _rv3 = str(d.get('投与量数値', 0) or 0).strip()
                _rv3 = ''.join(c for c in _rv3 if c.isdigit() or c == '.')
                dose_num = float(_rv3 or 0)
            except:
                dose_num = 0
            dose_str, unit_str = format_dose_text(d)
            if kubun == '抗がん剤':
                if dose_base=='BSA依存':
                    _dn = int(dose_num) if dose_num == int(dose_num) else dose_num
                    dose_parts.append(f'({_dn}mg/m²)')
                elif dose_base=='AUC依存': dose_parts.append(f'(AUC：{int(dose_num)})')
                elif dose_base=='BW依存':
                    _dn = int(dose_num) if dose_num == int(dose_num) else dose_num
                    dose_parts.append(f'({_dn}mg/kg)')
                elif dose_base=='固定用量':dose_parts.append(f'({dose_str}{unit_str})')
        time_val  = rp_sorted[0].get('投与時間文字','')
        time_text = f'(投与時間：{time_val})' if time_val else ''
        type_str  = ''
        for d in rp_sorted:
            t = str(master_dict.get(str(d.get('管理コード','')),{}).get('スケジュールシール用種類',''))
            if t: type_str = t; break
        drug_days = parse_days_num(rp_sorted[0].get('投与Day数値',''))
        return name_text, dose_parts, time_text, type_str, drug_days

    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mg=Cm(0.8); cw=A4_W-mg*2; x0=mg; y0=mg
    title_w=cw*0.72; date_w=cw*0.28
    title_h=Cm(1.5); course_h=Cm(1.0)

    add_textbox(slide,x0,y0,title_w,title_h,
                text=basic['レジメン名'],
                font_size=Pt(24),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0,y0+title_h,title_w,course_h,
                text=f'（1コース{cycle}日）',
                font_size=Pt(20),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0+title_w,y0,date_w,title_h+course_h,
                text='当レジメンの開始日\n年　　月　　日',
                font_size=Pt(11),bold=False,color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    y_cur = y0+title_h+course_h+Cm(0.2)
    note_lines = (['・'+str(n['注意事項文章']) for n in seal_notes] if seal_notes else [
        '・薬剤の他、副作用を予防する、投与ラインを満たす、薬剤を洗い流すために生理食塩液等を投与します。',
        '・強い副作用が出た場合や患者様の状態によっては、治療を延期したり、減量することがあります。'
    ])
    note_h = Cm(0.55*len(note_lines))
    add_textbox_multi(slide,x0,y_cur,cw,note_h,
                      lines=note_lines,font_size=Pt(11),
                      color=COLOR_BLACK,align=PP_ALIGN.LEFT)
    y_cur += note_h+Cm(0.2)

    has_right   = bool(oral_schedule) or need_bw or need_bsa
    tbl_w_ratio = 0.60 if has_right else 1.0
    right_w     = cw*0.38
    right_x     = x0+cw*tbl_w_ratio+cw*0.02

    type_w_cm=2.2; name_w_cm=3.8; inv_w_cm=1.1; rest_w_cm=1.5
    col_w_cm = [inv_w_cm if c['type']=='投与'
                else inv_w_cm*0.8 if c['type']=='休薬中'
                else rest_w_cm for c in columns]
    tbl_total_w  = type_w_cm+name_w_cm+sum(col_w_cm)
    scale        = float((cw*tbl_w_ratio)/Cm(tbl_total_w))
    type_w_cm   *= scale; name_w_cm *= scale
    col_w_cm     = [w*scale for w in col_w_cm]

    n_cols=2+len(columns); n_rps=len(sorted_rps); n_rows=1+n_rps

    def calc_lines(nd,dp,tt):
        return nd.count('\n')+1+len(dp)+(1 if tt else 0)

    row_heights = [Cm(0.65)]
    for rp in sorted_rps:
        nd,dp,tt,_,_ = get_rp_info(rp_groups[rp])
        row_heights.append(max(Cm(0.75),Cm(0.42*calc_lines(nd,dp,tt)+0.2)))
    tbl_h = sum(row_heights)

    shape = slide.shapes.add_table(n_rows,n_cols,x0,y_cur,
                                   Cm(tbl_total_w*scale),tbl_h)
    table = shape.table
    table.columns[0].width = Cm(type_w_cm)
    table.columns[1].width = Cm(name_w_cm)
    for i,w in enumerate(col_w_cm): table.columns[2+i].width = Cm(w)
    for i,h in enumerate(row_heights): table.rows[i].height = h

    for col_i,label in enumerate(['種類','薬剤名']):
        cell = table.cell(0,col_i)
        set_cell_text_multi(cell,[label],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)
    for i,col in enumerate(columns):
        cell = table.cell(0,2+i)
        set_cell_text_multi(cell,[col['label']],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)

    rest_col_indices = [(i,col) for i,col in enumerate(columns)
                        if col['type'] in ('休薬','休薬中')]
    for row_i,rp in enumerate(sorted_rps):
        r=row_i+1
        name_text,dose_parts,time_text,type_str,drug_days = get_rp_info(rp_groups[rp])
        row_color = TYPE_COLOR_MAP.get(type_str,COLOR_NAUSEA)
        cell = table.cell(r,0)
        set_cell_text_multi(cell,[type_str],font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        name_lines = name_text.split('\n')+dose_parts
        if time_text: name_lines.append(time_text)
        cell = table.cell(r,1)
        set_cell_text_multi(cell,name_lines,font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        for i,col in enumerate(columns):
            cell = table.cell(r,2+i); set_cell_margin(cell)
            if col['type']=='投与':
                hit = any(d in drug_days for d in col['days'])
                if hit:
                    set_cell_text_multi(cell,['●'],font_size=Pt(14),bold=True,
                                         color=COLOR_BLACK,align=PP_ALIGN.CENTER)
                    set_cell_bg(cell,row_color)
                else:
                    set_cell_bg(cell,COLOR_REST_BG)

    for i,col in rest_col_indices:
        col_idx = 2+i
        merge_cells_vertical(table,col_idx,1,n_rps)
        cell = table.cell(1,col_idx)
        set_cell_bg(cell,COLOR_REST_BG)
        set_cell_text_multi(cell,[''],font_size=Pt(11),
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_margin(cell)

    tbl_bottom  = y_cur+tbl_h
    y_after_tbl = tbl_bottom+Cm(0.15)
    if tegaki_notes:
        nb_h = Cm(0.5*len(tegaki_notes))
        add_textbox_multi(slide,x0,y_after_tbl,cw*tbl_w_ratio,nb_h,
                          lines=[str(n['注意事項文章']) for n in tegaki_notes],
                          font_size=Pt(11),color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    if has_right:
        ry = y_cur
        if oral_schedule:
            oral_lines = []
            for drug in oral_schedule:
                name   = str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
                timing = str(drug.get('投与タイミング',''))
                line   = f'・{name}'
                if timing and timing not in ('','ー','nan'):
                    line += f'\n　{timing}'
                oral_lines.append(line)
            oral_box_h = Cm(0.55*sum(l.count('\n')+1 for l in oral_lines)+0.3)
            add_textbox_multi(slide,right_x+Cm(0.15),ry+Cm(0.1),
                              right_w-Cm(0.3),oral_box_h,
                              lines=oral_lines,font_size=Pt(11),
                              color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += oral_box_h+Cm(0.3)
        if need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(0.7),
                        text='体重　　　　　kg',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(0.9)
        if need_bsa or need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(1.0),
                        text='体表面積\n\n＿＿＿＿m²',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(1.3)
        add_textbox(slide,right_x,ry,right_w,Cm(0.8),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        ry += Cm(1.0)
        shorten_notes = []
        for _d in drugs:
            _sn = str(_d.get('短縮注記','')).strip()
            if _sn and _sn not in shorten_notes:
                shorten_notes.append(_sn)
        if shorten_notes:
            for sn in shorten_notes:
                add_textbox(slide,right_x,ry,right_w,Cm(0.6),
                            text=sn,font_size=Pt(9),bold=False,
                            color=RGBColor(0x85,0x64,0x04),
                            align=PP_ALIGN.LEFT)
                ry += Cm(0.65)
        add_textbox(slide,right_x,ry,right_w,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
    else:
        footer_y = A4_H-Cm(1.5)-mg
        add_textbox(slide,x0,footer_y,cw*0.5,Cm(0.7),
                    text='体表面積　　　　m²',font_size=Pt(11),
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        add_textbox(slide,x0+cw*0.5,footer_y,cw*0.5,Cm(0.7),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
        shorten_notes_b = []
        for _d in drugs:
            _sn = str(_d.get('短縮注記','')).strip()
            if _sn and _sn not in shorten_notes_b:
                shorten_notes_b.append(_sn)
        if shorten_notes_b:
            for sn in shorten_notes_b:
                add_textbox(slide,x0,footer_y+Cm(0.8),cw,Cm(0.5),
                            text=sn,font_size=Pt(9),bold=False,
                            color=RGBColor(0x85,0x64,0x04),
                            align=PP_ALIGN.LEFT)
                footer_y += Cm(0.55)
        add_textbox(slide,x0,footer_y+Cm(0.8),cw,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ===== STEP1 =====
st.subheader("STEP 1　レジメンPDFをアップロード")
uploaded = st.file_uploader(
    "PDFファイルをここにドロップ",
    type="pdf",
    help="レジメン情報PDFを1件アップロードしてください"
)
if uploaded:
    st.success(f"✅ {uploaded.name} を読み込みました")
st.divider()

# ===== STEP2: AI自動解析 =====
st.subheader("STEP 2　AIが自動解析")
if uploaded:
    if st.button("🤖 自動解析スタート", type="primary", use_container_width=True):
        with st.spinner("AIが解析中です...少々お待ちください⏳"):
            try:
                definition = load_definition()
                pdf_bytes  = uploaded.read()
                client     = get_gemini_client()
                response   = client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[
                        definition,
                        types.Part.from_bytes(mime_type="application/pdf", data=pdf_bytes)
                    ]
                )
                raw   = response.text
                match = re.search(r"\{.*\}", raw, re.DOTALL)
                if match:
                    json_str = match.group()
                    parsed   = json.loads(json_str)
                    st.session_state["extracted_json"]   = json_str
                    st.session_state["extracted_parsed"] = parsed
                    st.session_state["json_editor_sync"] = True
                    st.session_state.pop("registered", None)
                    st.session_state.pop("yonin_confirmed_1", None)
                    # STEP3.5のキーもリセット
                    for _k in list(st.session_state.keys()):
                        if _k.startswith("step35_"):
                            del st.session_state[_k]
                    st.success("✅ 解析完了！")
                else:
                    st.error("JSONの抽出に失敗しました。もう一度試してください。")
                    st.text(raw)
            except Exception as e:
                st.error(f"エラーが発生しました: {e}")
else:
    st.info("👆 まずPDFをアップロードしてください")
st.divider()

# ===== STEP3: 内容確認 =====
st.subheader("STEP 3　内容を確認・修正")
if "extracted_parsed" in st.session_state:
    parsed = st.session_state["extracted_parsed"]
    if "pattern_determination" in parsed:
        st.info(parsed["pattern_determination"])
    basic_json = parsed.get("basic_info", {})
    drug_list  = parsed.get("drug_info", [])

    st.markdown("#### 📋 基本情報")
    col1, col2 = st.columns(2)
    with col1:
        st.write(f"**プロトコールNo：** {basic_json.get('protocol_no','')}")
        st.write(f"**レジメン名：** {basic_json.get('regimen_name','')}")
        st.write(f"**対象疾患：** {basic_json.get('disease','')}")
        st.write(f"**疾患分類：** {basic_json.get('disease_category','')}")
    with col2:
        st.write(f"**1コース日数：** {basic_json.get('course_days','')}")
        st.write(f"**備考：** {basic_json.get('remarks','')}")
        st.write(f"**登録日：** {basic_json.get('registration_date','')}")

    st.markdown("#### 💊 薬剤情報")
    if drug_list:
        import pandas as pd
        st.dataframe(pd.DataFrame(drug_list), use_container_width=True)

    with st.expander("🔧 JSONを直接編集する"):
        with st.form(key="form_json_editor"):
            _edited = st.text_area(
                "JSON",
                value=st.session_state.get("extracted_json", ""),
                height=400,
            )
            _json_submitted = st.form_submit_button(
                "✅ 編集内容を反映",
                type="primary",
                use_container_width=True
            )
        if _json_submitted:
            try:
                _parsed_new = json.loads(_edited)
                st.session_state["extracted_json"]   = _edited
                st.session_state["extracted_parsed"] = _parsed_new
                st.session_state.pop("yonin_confirmed_1", None)
                st.session_state.pop("json_editor_text", None)
                # STEP3.5のキーもリセット
                for _k in list(st.session_state.keys()):
                    if _k.startswith("step35_"):
                        del st.session_state[_k]
                st.success("✅ 反映しました！")
                st.rerun()
            except json.JSONDecodeError as e:
                st.error(f"JSONの形式が正しくありません: {e}")
else:
    st.info("👆 STEP2で解析すると結果がここに表示されます")
st.divider()

# ===== STEP3.5: 要確認項目の解決 =====
st.subheader("STEP 3.5　要確認項目の入力")

if "extracted_parsed" in st.session_state and not st.session_state.get("registered"):
    import re as _re35

    def _parse_time_hours_35(text):
        text = str(text).strip()
        h = _re35.search(r"(\d+)\s*時間", text)
        m = _re35.search(r"(\d+)\s*分",   text)
        hours   = int(h.group(1)) if h else 0
        minutes = int(m.group(1)) if m else 0
        return hours + minutes / 60

    _parsed35 = st.session_state["extracted_parsed"]
    _drugs35  = _parsed35.get("drug_info") or _parsed35.get("drugs") or []
    _basic35  = _parsed35.get("basic_info", {})

    YONIN_DEF = [
        ("admin_time_text",   "投与時間",    "例：30分・1時間・2時間"),
        ("admin_day_text",    "投与Day",     "例：day1・day1,8,15"),
        ("admin_day_numeric", "投与Day数値", "例：1・1|8|15"),
        ("diluent_volume",    "希釈液容量",  "例：250"),
        ("management_code",   "管理コード",  "例：AC001"),
        ("dosage_value",      "投与量数値",  "例：100"),
    ]
    YONIN_BASIC = [
        ("course_days", "1コース日数",    "例：21・28"),
        ("protocol_no", "プロトコールNo", "例：C34-001"),
    ]

    _yonin_basic_items = []
    for _k, _label, _placeholder in YONIN_BASIC:
        _val = str(_basic35.get(_k, "")).strip()
        if _val == "要確認":
            _yonin_basic_items.append({
                "key"        : f"step35_basic_{_k}",
                "json_key"   : _k,
                "label"      : f"基本情報 / {_label}",
                "placeholder": _placeholder,
            })

    _yonin_drug_items = []
    for _di, _d in enumerate(_drugs35):
        _dname = str(
            _d.get("product_name") or _d.get("brand_name") or
            _d.get("management_code") or f"薬剤{_di+1}"
        ).strip()
        for _k, _label, _placeholder in YONIN_DEF:
            _val = str(_d.get(_k, "") or "").strip()
            if _val == "要確認":
                _yonin_drug_items.append({
                    "key"        : f"step35_drug_{_di}_{_k}",
                    "drug_idx"   : _di,
                    "drug_name"  : _dname,
                    "json_key"   : _k,
                    "label"      : f"{_dname} / {_label}",
                    "placeholder": _placeholder,
                })

    _all_yonin = _yonin_basic_items + _yonin_drug_items

    if not _all_yonin:
        st.success("✅ 要確認項目はありません。STEP4へ進んでください。")
    else:
        st.warning(f"⚠️ {len(_all_yonin)} 件の要確認項目があります。入力してから登録してください。")
        st.caption("未入力のままでも登録できますが、スプレッドシートに要確認が残ります。")

        _all_filled = True
        for _item in _all_yonin:
            _inp = st.text_input(
                _item["label"],
                placeholder=_item["placeholder"],
                key=_item["key"],
            )
            if not _inp.strip():
                _all_filled = False

        st.divider()
        col_fix, col_skip = st.columns(2)

        with col_fix:
            if st.button(
                "✅ 入力内容をJSONに反映してSTEP4へ",
                type="primary",
                use_container_width=True,
                key="btn_step35_fix",
                disabled=not _all_filled
            ):
                import copy
                _parsed_new = copy.deepcopy(_parsed35)
                _basic_new  = _parsed_new.get("basic_info", {})
                _drugs_new  = (
                    _parsed_new.get("drug_info") or
                    _parsed_new.get("drugs") or []
                )
                for _item in _yonin_basic_items:
                    _v = st.session_state.get(_item["key"], "").strip()
                    if _v:
                        _k = _item["json_key"]
                        if _k == "course_days":
                            try:
                                _basic_new[_k] = int(_v)
                            except:
                                _basic_new[_k] = _v
                        else:
                            _basic_new[_k] = _v
                for _item in _yonin_drug_items:
                    _v = st.session_state.get(_item["key"], "").strip()
                    if _v:
                        _di = _item["drug_idx"]
                        _k  = _item["json_key"]
                        if _di < len(_drugs_new):
                            if _k == "admin_time_text":
                                _drugs_new[_di][_k] = _v
                                _drugs_new[_di]["admin_time_numeric"] = round(
                                    _parse_time_hours_35(_v), 4
                                )
                            elif _k == "diluent_volume":
                                try:
                                    _drugs_new[_di][_k] = float(_v)
                                except:
                                    _drugs_new[_di][_k] = _v
                            else:
                                _drugs_new[_di][_k] = _v

                _parsed_new["basic_info"] = _basic_new
                if "drug_info" in _parsed_new:
                    _parsed_new["drug_info"] = _drugs_new
                else:
                    _parsed_new["drugs"] = _drugs_new

                st.session_state["extracted_parsed"] = _parsed_new
                st.session_state["extracted_json"]   = json.dumps(
                    _parsed_new, ensure_ascii=False, indent=2
                )
                for _item in _all_yonin:
                    st.session_state.pop(_item["key"], None)
                st.success("✅ JSONに反映しました。STEP4で登録してください。")
                st.rerun()

        with col_skip:
            if st.button(
                "⏭️ このままSTEP4へ（要確認を残す）",
                use_container_width=True,
                key="btn_step35_skip"
            ):
                st.rerun()

else:
    if not st.session_state.get("registered"):
        st.info("👆 STEP2で解析すると要確認項目が表示されます")
st.divider()

# ===== STEP4: 登録 =====
st.subheader("STEP 4　スプレッドシートに登録")
if "extracted_parsed" in st.session_state and not st.session_state.get("registered"):
    data         = st.session_state["extracted_parsed"]
    info         = get_basic(data)
    drugs        = get_drugs(data)
    protocol_no  = get_val(info, "protocol_no")
    regimen_name = get_val(info, "regimen_name")
    disease      = get_val(info, "disease", "target_disease", "対象疾患", default="")
    course_days  = get_val(info, "course_days", default="要確認")

    if not protocol_no or protocol_no == "要確認":
        st.warning("⚠️ プロトコールNoが「要確認」です。JSONを編集してから登録してください。")
    else:
        sh        = get_spreadsheet()
        ws_basic  = sh.worksheet("基本情報")
        existing  = ws_basic.col_values(1)
        is_dup    = protocol_no in existing

        if is_dup:
            st.warning(f"⚠️ {protocol_no} はすでに登録されています。上書きしますか？")
            col1, col2 = st.columns(2)
            with col1:
                if st.button("❌ キャンセル"):
                    st.session_state.pop("extracted_parsed", None)
                    st.session_state.pop("extracted_json", None)
                    st.rerun()
            with col2:
                do_register = st.button("⚠️ 上書き登録", type="primary")
            overwrite = True
        else:
            st.success(f"✅ {protocol_no} は新規登録です。")
            do_register = st.button("✅ 新規登録", type="primary", use_container_width=True)
            overwrite = False

        if do_register:
            with st.spinner("登録中..."):
                try:
                    today     = date.today().strftime("%Y/%-m/%-d")
                    now       = datetime.now().strftime("%Y/%-m/%-d %H:%M")
                    ws_drug   = sh.worksheet("薬剤情報")
                    ws_log    = sh.worksheet("抽出ログ")

                    if overwrite:
                        existing2 = ws_basic.col_values(1)
                        if protocol_no in existing2:
                            ws_basic.delete_rows(existing2.index(protocol_no)+1)
                        all_vals = ws_drug.get_all_values()
                        for i in reversed([i+1 for i,r in enumerate(all_vals)
                                           if r and r[0]==protocol_no]):
                            ws_drug.delete_rows(i)

                    basic_row = [
                        protocol_no, regimen_name, disease, "",
                        course_days if course_days else "要確認",
                        "","","","", today, "",
                    ]
                    ws_basic.append_row(basic_row, value_input_option="USER_ENTERED")

                    for drug in drugs:
                        drug_row = [
                            protocol_no,
                            get_val(drug,"order"),
                            get_val(drug,"management_code"),
                            get_val(drug,"product_name","brand_name"),
                            get_val(drug,"dosage_value","dose_value"),
                            get_val(drug,"dosage_unit","dose_unit"),
                            get_val(drug,"dosage_basis","dose_basis"),
                            get_val(drug,"admin_day_text","day_text"),
                            get_val(drug,"admin_day_numeric","day_numbers"),
                            get_val(drug,"admin_timing","timing"),
                            get_val(drug,"diluent_volume","diluent_ml"),
                            get_val(drug,"admin_time_text","infusion_time_text"),
                            get_val(drug,"admin_time_numeric","infusion_time_hours"),
                            get_val(drug,"anticancer_flag","flag_O_chemo"),
                            get_val(drug,"support_flag","flag_O_support"),
                            get_val(drug,"seal_flag","flag_seal"),
                            get_val(drug,"figure_flag","flag_chart"),
                            get_val(drug,"manual_flag","flag_leaflet"),
                            get_val(drug,"remarks","note"),
                        ]
                        ws_drug.append_row(drug_row, value_input_option="USER_ENTERED")

                    # Pdカテゴリ自動設定
                    try:
                        ws_ae    = sh.worksheet("抗がん剤副作用マスタ")
                        ws_pd_sh = sh.worksheet("Pd")
                        ae_data  = ws_ae.get_all_records()
                        pd_data  = ws_pd_sh.get_all_records()
                        trigger_to_pdid = {}
                        code_to_pdid    = {}
                        priority_dict   = {}
                        for p in pd_data:
                            trigger = str(p.get('トリガーキーワード','')).strip()
                            cat_id  = str(p.get('カテゴリID','')).strip()
                            try: priority_dict[cat_id] = int(p.get('優先順位',99))
                            except: priority_dict[cat_id] = 99
                            if not trigger or trigger == '手動設定': continue
                            if trigger.startswith('AC'):
                                for c in trigger.split('|'):
                                    code_to_pdid[c.strip()] = cat_id
                            else:
                                if trigger not in trigger_to_pdid:
                                    trigger_to_pdid[trigger] = []
                                trigger_to_pdid[trigger].append(cat_id)
                        ae_dict    = {str(r['管理コード']).strip(): r for r in ae_data}
                        ae_headers = ws_ae.row_values(1)
                        ae_columns = ae_headers[2:]
                        collected  = set()
                        for drug in drugs:
                            drug_code = str(get_val(drug,'management_code')).strip()
                            ae_row    = ae_dict.get(drug_code)
                            if ae_row:
                                for col in ae_columns:
                                    if str(ae_row.get(col,'')).strip() == '○':
                                        for pid in trigger_to_pdid.get(col,[]):
                                            collected.add(pid)
                            if drug_code in code_to_pdid:
                                collected.add(code_to_pdid[drug_code])
                        sorted_ids  = sorted(collected, key=lambda x: priority_dict.get(x,99))
                        pd_category = '|'.join(sorted_ids)
                        if pd_category:
                            existing3 = ws_basic.col_values(1)
                            if protocol_no in existing3:
                                ws_basic.update_cell(existing3.index(protocol_no)+1, 12, pd_category)
                    except Exception as pd_e:
                        st.warning(f"⚠️ Pdカテゴリ自動設定エラー: {pd_e}")

                    ws_log.append_row([
                        now, protocol_no, regimen_name,
                        "上書き登録" if overwrite else "新規登録", len(drugs)
                    ], value_input_option="USER_ENTERED")

                    st.session_state["registered"]          = True
                    st.session_state["registered_protocol"] = protocol_no
                    st.success(f"✅ {protocol_no} を登録しました！")
                    st.balloons()
                    st.rerun()

                except Exception as e:
                    st.error(f"❌ 登録エラー: {e}")

elif st.session_state.get("registered"):
    st.success(f"✅ {st.session_state.get('registered_protocol','')} 登録済み")
else:
    st.info("👆 STEP3で内容を確認してから登録してください")

st.divider()

# ===== STEP5: パワポ生成 =====
st.subheader("STEP 5　（外来化療）スケジュールシール パワポ生成")

if st.session_state.get("registered"):
    protocol_no_reg = st.session_state.get("registered_protocol", "")
    st.info(f"📋 対象レジメン：{protocol_no_reg}")

    if st.button("📑 パワポ（外来スケジュールシール）作成",
                 type="primary", use_container_width=True):
        with st.spinner("パワポ生成中..."):
            try:
                sh          = get_spreadsheet()
                basic_data  = sh.worksheet("基本情報").get_all_records()
                drug_data   = sh.worksheet("薬剤情報").get_all_records()
                master_data, notes_data = load_master_data()
                pptx_data   = create_pptx(
                    protocol_no_reg, basic_data, drug_data, master_data, notes_data
                )
            except Exception as e:
                st.error(f"データ取得エラー: {e}")
                pptx_data = None

        if pptx_data:
            st.download_button(
                label="⬇️ パワポをダウンロード",
                data=pptx_data,
                file_name=f"{protocol_no_reg}_スケジュールシール_AI作成要チェック_{today_str}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
            st.success("✅ パワポ生成完了！")
        else:
            st.error("❌ 生成に失敗しました")

    st.divider()
    st.info("続けてどうしますか？")

    if st.button(
        "📋 パラメータ入力しテキスト生成ページへ（O欄・Pd欄テキストコピー）",
        use_container_width=True,
        key="btn_next_text"
    ):
        st.switch_page("pages/4_O欄Pd欄生成.py")

    if st.button(
        "📊 テンプレート生成ページへ（Excel生成）",
        use_container_width=True,
        key="btn_next_excel"
    ):
        st.switch_page("pages/3_テンプレート生成.py")

    if st.button(
        "🔄 続けて別のレジメンを登録する",
        use_container_width=True,
        key="btn_next_new"
    ):
        for key in [
            "extracted_json", "extracted_parsed",
            "registered", "registered_protocol",
            "yonin_confirmed_1",
            "json_editor_text", "json_editor_sync",
        ]:
            st.session_state.pop(key, None)
        for key in list(st.session_state.keys()):
            if key.startswith("step35_") or key.startswith("yonin1_"):
                del st.session_state[key]
        st.rerun()
else:
    st.info("👆 STEP4で登録が完了するとパワポ生成ボタンが表示されます")
