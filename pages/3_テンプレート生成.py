import streamlit as st
import re
import io
import gspread
from google.oauth2 import service_account
from collections import defaultdict
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import urllib.request
import csv
from datetime import date

today_str = date.today().strftime("%Y%m%d")

st.set_page_config(
    page_title="テンプレート生成",
    page_icon="💊",
    layout="centered",
    initial_sidebar_state="expanded",
    menu_items={}
)

st.sidebar.title("メニュー")

SPREADSHEET_URL = "https://docs.google.com/spreadsheets/d/1dLEUYSZlrIK1uHqEtEAfS1jSAPpXCIiAiAk_iaRuY-8/edit"

@st.cache_resource
def get_gspread_client():
    creds = service_account.Credentials.from_service_account_info(
        st.secrets["gcp_service_account"],
        scopes=[
            "https://spreadsheets.google.com/feeds",
            "https://www.googleapis.com/auth/drive",
        ]
    )
    return gspread.authorize(creds)

@st.cache_data(ttl=60)
def fetch_sheet(sheet_name):
    try:
        gc = get_gspread_client()
        ss = gc.open_by_url(SPREADSHEET_URL)
        ws = ss.worksheet(sheet_name)
        return ws.get_all_records()
    except Exception as e:
        st.error(f"シート「{sheet_name}」の取得に失敗: {e}")
        return []

@st.cache_data(ttl=300)
def load_all_data():
    basic_data  = fetch_sheet("基本情報")
    drug_data   = fetch_sheet("薬剤情報")
    master_data = fetch_sheet("薬品マスタ")
    notes_data  = fetch_sheet("注意事項")
    pd_data     = fetch_sheet("Pd")
    # ===== 追加：抗がん剤副作用マスタを読み込む =====
    ae_data     = fetch_sheet("抗がん剤副作用マスタ")
    return basic_data, drug_data, master_data, notes_data, pd_data, ae_data

def to_half_kana(text):
    table = {
        'ア':'ｱ','イ':'ｲ','ウ':'ｳ','エ':'ｴ','オ':'ｵ',
        'カ':'ｶ','キ':'ｷ','ク':'ｸ','ケ':'ｹ','コ':'ｺ',
        'サ':'ｻ','シ':'ｼ','ス':'ｽ','セ':'ｾ','ソ':'ｿ',
        'タ':'ﾀ','チ':'ﾁ','ツ':'ﾂ','テ':'ﾃ','ト':'ﾄ',
        'ナ':'ﾅ','ニ':'ﾆ','ヌ':'ﾇ','ネ':'ﾈ','ノ':'ﾉ',
        'ハ':'ﾊ','ヒ':'ﾋ','フ':'ﾌ','ヘ':'ﾍ','ホ':'ﾎ',
        'マ':'ﾏ','ミ':'ﾐ','ム':'ﾑ','メ':'ﾒ','モ':'ﾓ',
        'ヤ':'ﾔ','ユ':'ﾕ','ヨ':'ﾖ',
        'ラ':'ﾗ','リ':'ﾘ','ル':'ﾙ','レ':'ﾚ','ロ':'ﾛ',
        'ワ':'ﾜ','ヲ':'ｦ','ン':'ﾝ',
        'ァ':'ｧ','ィ':'ｨ','ゥ':'ｩ','ェ':'ｪ','ォ':'ｫ',
        'ッ':'ｯ','ャ':'ｬ','ュ':'ｭ','ョ':'ｮ',
        'ガ':'ｶﾞ','ギ':'ｷﾞ','グ':'ｸﾞ','ゲ':'ｹﾞ','ゴ':'ｺﾞ',
        'ザ':'ｻﾞ','ジ':'ｼﾞ','ズ':'ｽﾞ','ゼ':'ｾﾞ','ゾ':'ｿﾞ',
        'ダ':'ﾀﾞ','ヂ':'ﾁﾞ','ヅ':'ﾂﾞ','デ':'ﾃﾞ','ド':'ﾄﾞ',
        'バ':'ﾊﾞ','ビ':'ﾋﾞ','ブ':'ﾌﾞ','ベ':'ﾍﾞ','ボ':'ﾎﾞ',
        'パ':'ﾊﾟ','ピ':'ﾋﾟ','プ':'ﾌﾟ','ペ':'ﾍﾟ','ポ':'ﾎﾟ',
        'ー':'ｰ','ヴ':'ｳﾞ','・':'･',
    }
    result = ''
    for char in str(text):
        result += table.get(char, char)
    return result

def shorten_regimen_name(regimen_name):
    name = regimen_name
    name = re.sub(r'[（(][^）)]*[）)]', '', name)
    for word in [
        '術前','術後','周術期','切除不能','再発','進行',
        '維持','補助','一次治療','二次治療','初回','難治性',
        '肺癌','胃癌','大腸癌','乳癌','膵癌','肝癌',
        '食道癌','子宮癌','卵巣癌','前立腺癌','膀胱癌',
        '腎癌','甲状腺癌','悪性リンパ腫','白血病',
        '骨髄腫','中皮腫','胸腺癌','胸腺腫','神経内分泌腫瘍',
    ]:
        name = name.replace(word, '')
    name = re.sub(r'\s+', ' ', name).strip('　 ')
    return name

def parse_days_num(day_str):
    days = []
    for part in str(day_str).split('|'):
        part = part.strip()
        if '-' in part:
            try:
                s, e = part.split('-')
                days.extend(range(int(s), int(e)+1))
            except:
                pass
        elif part.isdigit():
            days.append(int(part))
    return sorted(set(days))

def format_dose_text(drug):
    try:
        dose = float(drug['投与量数値'])
    except:
        dose = 0
    unit_input = str(drug.get('投与単位', '')).strip()
    dose_base  = str(drug.get('用量根拠', ''))
    v_to_mg    = drug.get('1V当たりmg', '')
    if unit_input.upper() == 'V':
        if v_to_mg != '' and str(v_to_mg).strip() != '':
            try:
                dose = dose * float(v_to_mg)
                unit_input = 'mg'
            except:
                pass
    if unit_input in ('mg','mg/body','mg/ body'):
        display_unit = 'mg'
    elif unit_input == '':
        display_unit = '' if dose_base == 'AUC依存' else 'mg'
    else:
        display_unit = unit_input
    dose_str = str(int(dose)) if dose == int(dose) else str(dose)
    return dose_str, display_unit

INJECTION_ORDER = {
    'NK1':1,'5HT3':2,'ステロイド':3,'G-CSF':4,
    '利尿薬':5,'解毒薬':6,'抗アレルギー':7,
    'H2ブロッカー':8,'電解質補正':9,'その他注射':10
}

def make_support_line(drug):
    name = to_half_kana(
        str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
    )
    dose_str, unit_str = format_dose_text(drug)
    day = str(drug.get('投与Day文字',''))
    return f"{name} {dose_str}{unit_str}({day})"

def get_regimen(protocol_no, basic_data, drug_data, master_data):
    master_dict = {m['管理コード']: m for m in master_data}
    basic = [b for b in basic_data
             if b['プロトコールNo'] == protocol_no]
    if not basic:
        return None
    basic = basic[0]
    drugs_raw = sorted(
        [d for d in drug_data
         if d['プロトコールNo'] == protocol_no],
        key=lambda x: (
            0 if str(x['投与順序']).isdigit() else 1,
            int(x['投与順序']) if str(x['投与順序']).isdigit() else 99
        )
    )
    drugs = []
    for drug in drugs_raw:
        code   = str(drug['管理コード'])
        master = master_dict.get(code, {})
        merged = dict(drug)
        merged.update({
            '一般名（全角）'        : master.get('一般名（全角）',''),
            '一般名（半角カナ）'    : master.get('一般名（半角カナ）',''),
            '採用商品名（全角）'    : master.get('採用商品名（全角）',''),
            '採用商品名（半角カナ）': master.get('採用商品名（半角カナ）',''),
            '薬効分類'             : master.get('薬効分類',''),
            '薬剤区分'             : master.get('薬剤区分',''),
            '支持療法分類'         : master.get('支持療法分類',''),
            '薬品マスタ単位'       : master.get('単位',''),
            '投与経路'             : master.get('投与経路',''),
            '1V当たりmg'           : master.get('1V当たりmg',''),
            '患者向け説明'         : master.get('患者向け説明',''),
        })
        drugs.append(merged)
    return {'basic': basic, 'drugs': drugs, 'master_dict': master_dict}



# ===== 追加：副作用登録UI関数 =====
def show_ae_register_ui(unregistered, ae_data, master_data, drug_data, basic_data, pd_data):
    """未登録薬剤の副作用をその場で登録するUI"""

    # 薬品マスタを辞書化（管理コード→レコード）
    master_dict = {str(m['管理コード']).strip(): m for m in master_data}

    # 副作用列名リスト
    gc_client   = get_gspread_client()
    ss          = gc_client.open_by_url(SPREADSHEET_URL)
    ws_ae       = ss.worksheet("抗がん剤副作用マスタ")
    ws_basic    = ss.worksheet("基本情報")
    ae_headers  = ws_ae.row_values(1)
    ae_columns  = ae_headers[2:-1]  # 登録日を除く副作用列
    ae_all      = ws_ae.get_all_values()

    # セッション初期化
    if "ae_reg_index" not in st.session_state:
        st.session_state["ae_reg_index"] = 0
    if "ae_reg_done" not in st.session_state:
        st.session_state["ae_reg_done"] = []

    idx = st.session_state["ae_reg_index"]

    # 全部完了した場合
    if idx >= len(unregistered):
        st.success("✅ 全薬剤の副作用登録が完了しました！")

        # Pdカテゴリ自動更新
        with st.spinner("Pdカテゴリを自動更新中..."):
            try:
                ws_pd_sh    = ss.worksheet("Pd")
                pd_data_new = ws_pd_sh.get_all_records()
                ae_data_new = ws_ae.get_all_records()

                trigger_to_pdid = {}
                code_to_pdid    = {}
                priority_dict   = {}
                for p in pd_data_new:
                    trigger = str(p.get('トリガーキーワード', '')).strip()
                    cat_id  = str(p.get('カテゴリID', '')).strip()
                    try:
                        priority_dict[cat_id] = int(p.get('優先順位', 99))
                    except:
                        priority_dict[cat_id] = 99
                    if not trigger or trigger == '手動設定':
                        continue
                    if trigger.startswith('AC'):
                        for c in trigger.split('|'):
                            code_to_pdid[c.strip()] = cat_id
                    else:
                        trigger_to_pdid[trigger] = cat_id

                ae_dict_new = {str(r['管理コード']).strip(): r for r in ae_data_new}
                ae_cols_new = ws_ae.row_values(1)[2:-1]

                # このレジメンのPdカテゴリを再計算
                protocol_no = st.session_state.get("current_protocol_no", "")
                if protocol_no:
                    reg_drugs = [
                        d for d in drug_data
                        if str(d.get('プロトコールNo', '')).strip() == protocol_no
                    ]
                    collected = set()
                    for drug in reg_drugs:
                        drug_code = str(drug.get('管理コード', '')).strip()
                        ae_row = ae_dict_new.get(drug_code)
                        if ae_row:
                            for col in ae_cols_new:
                                if str(ae_row.get(col, '')).strip() == '○':
                                    pd_id = trigger_to_pdid.get(col)
                                    if pd_id:
                                        collected.add(pd_id)
                        if drug_code in code_to_pdid:
                            collected.add(code_to_pdid[drug_code])

                    sorted_ids  = sorted(collected, key=lambda x: priority_dict.get(x, 99))
                    pd_category = '|'.join(sorted_ids)

                    if pd_category:
                        basic_codes = ws_basic.col_values(1)
                        if protocol_no in basic_codes:
                            row_idx = basic_codes.index(protocol_no) + 1
                            ws_basic.update_cell(row_idx, 12, pd_category)
                            st.success(f"✅ Pdカテゴリを更新しました：{pd_category}")

                # キャッシュクリア
                st.cache_data.clear()
                st.session_state.pop("ae_reg_index", None)
                st.session_state.pop("ae_reg_done", None)
                st.session_state.pop("unregistered_drugs", None)

            except Exception as e:
                st.warning(f"⚠️ Pdカテゴリ更新エラー: {e}")

        st.info("ファイル生成に進んでください")
        return

    # 現在登録する薬剤
    drug_info = unregistered[idx]
    code      = drug_info['code']
    name      = drug_info['name']

    st.divider()
    st.subheader(f"💊 副作用登録 ({idx+1}/{len(unregistered)}）：{name}")

    # ---------- NotebookLM指示文 ----------
    with st.expander("📋 NotebookLMで抽出する場合はこちら", expanded=False):
        NOTEBOOKLM_URL = "https://notebooklm.google.com/"
        # ダウンロード強制リンク（?raw=true でファイルとして取得）
        DEFINITION_URL = (
            "https://github.com/chemo-pharm-tohoku/chemo-template-app"
            "/raw/main/%E5%89%AF%E4%BD%9C%E7%94%A8%E6%8A%BD%E5%87%BA"
            "%E5%AE%9A%E7%BE%A9%E6%9B%B8260720.txt"
        )
        # 薬品マスタから採用商品名（全角）を取得
        brand_name = next(
            (m.get('採用商品名（全角）', name)
             for m in master_data
             if str(m.get('管理コード', '')).strip() == code),
            name
        )
        brand_display = f"{name}（{brand_name}）" if brand_name != name else name

        st.markdown(
            f"**① [NotebookLMを開く]({NOTEBOOKLM_URL})** に以下をアップロードしてください"
        )
        st.info(
            f"📄 **副作用抽出定義書** → "
            f"[こちらからダウンロード]({DEFINITION_URL})\n\n"
            f"📄 **{brand_display} の添付文書PDF** → "
            f"[PMDAで検索](https://www.pmda.go.jp/PmdaSearch/iyakuSearch/) "
            f"で「{brand_name}」を検索してダウンロード"
        )
        st.divider()
        st.markdown("**② 以下の指示文をNotebookLMに貼り付けてください**")
        instruction = (
            f"副作用抽出定義書に従い、{name}の添付文書から\n"
            f"副作用情報を抽出しCSV形式で出力してください。\n"
            f"管理コードは「{code}」を使用してください。"
        )
        st.code(instruction, language="text")
        st.divider()
        st.markdown("**③ 出力されたCSVをここに貼り付けてください**")

        csv_input = st.text_input(
            "CSVを貼り付け（例：AC999,イピリムマブ,○,○,,,,,○）",
            key=f"csv_input_{code}",
            placeholder=f"{code},{name},○,○,,,,,"
        )

        if csv_input.strip():
            # CSVをパースしてチェックボックスに反映
            try:
                parts = [p.strip() for p in csv_input.split(',')]
                if len(parts) >= len(ae_columns) + 2:
                    # セッションに保存してrerunで反映
                    changed = False
                    for i, col in enumerate(ae_columns):
                        val     = parts[i + 2] if i + 2 < len(parts) else ''
                        new_val = (val == '○')
                        key     = f"cb_{code}_{col}"
                        if st.session_state.get(key) != new_val:
                            st.session_state[key] = new_val
                            changed = True
                    if changed:
                        st.session_state[f"csv_applied_{code}"] = True
                        st.rerun()
                    else:
                        pass  # 既に反映済み
                else:
                    st.warning(f"⚠️ CSV列数が不足しています（{len(parts)}列、必要：{len(ae_columns)+2}列）")
            except Exception as e:
                st.warning(f"⚠️ CSV解析エラー: {e}")

        # CSV反映完了メッセージは削除

    # ---------- チェックボックス ----------
    st.markdown("**副作用をチェックしてください**")

    # スプレッドシートの現在値を取得してデフォルト表示
    ae_dict_current = {str(r.get("管理コード","")).strip(): r for r in ae_data}
    current_ae_row  = ae_dict_current.get(code, {})

    cols = st.columns(3)
    checked = {}
    for i, col_name in enumerate(ae_columns):
        cb_key = f"cb_{code}_{col_name}"
        ae_key = f"ae_{code}_{col_name}"
        # CSV反映時のキー移行
        if ae_key in st.session_state and cb_key not in st.session_state:
            st.session_state[cb_key] = st.session_state.pop(ae_key)
        # セッションになければスプレッドシートの現在値をデフォルトに
        if cb_key not in st.session_state:
            current_val = str(current_ae_row.get(col_name, "")).strip()
            st.session_state[cb_key] = (current_val == "○")
        default = st.session_state.get(cb_key, False)
        with cols[i % 3]:
            checked[col_name] = st.checkbox(
                col_name,
                value=default,
                key=cb_key
            )

    # チェックリセットボタン
    if st.button("チェックをリセット", key=f"btn_reset_{code}"):
        for col_name in ae_columns:
            st.session_state.pop(f"cb_{code}_{col_name}", None)
        st.session_state.pop(f"csv_applied_{code}", None)
        st.rerun()

    # ---------- 登録・スキップ・リロードボタン ----------
    st.divider()
    col_reg, col_skip, col_reload = st.columns(3)

    with col_reg:
        if st.button(
            f"✅ {name} の副作用を登録する",
            type="primary",
            use_container_width=True,
            key=f"btn_reg_{code}"
        ):
            try:
                from datetime import date
                today = date.today().strftime("%Y/%m/%d")

                # 副作用マスタの行を探して更新
                ae_codes = [row[0] for row in ae_all]
                if code in ae_codes:
                    row_idx = ae_codes.index(code) + 1
                    # セッションステートから最新のチェック状態を取得
                    final_checked = {
                        col: st.session_state.get(f"cb_{code}_{col}", False)
                        for col in ae_columns
                    }
                    update_vals = [
                        ['○' if final_checked.get(col, False) else '' for col in ae_columns]
                        + [today]
                    ]
                    # 列文字を正確に指定
                    from openpyxl.utils import get_column_letter as gcl
                    start_col = gcl(3)
                    end_col   = gcl(2 + len(ae_columns) + 1)
                    ws_ae.update(
                        range_name=f'{start_col}{row_idx}:{end_col}{row_idx}',
                        values=update_vals
                    )
                    # デバッグ用（確認後削除可）
                    st.write(f"更新値: {update_vals}")
                    st.success(f"✅ {name} の副作用を登録しました！")
                    st.session_state["ae_reg_done"].append(code)
                    st.session_state["ae_reg_index"] += 1
                    st.rerun()
                else:
                    st.error(f"❌ {code} が副作用マスタに見つかりません")
            except Exception as e:
                st.error(f"❌ 登録エラー: {e}")

    with col_skip:
        if st.button(
            f"⏭️ {name} をスキップ",
            use_container_width=True,
            key=f"btn_skip_{code}"
        ):
            st.session_state["ae_reg_index"] += 1
            st.rerun()

    with col_reload:
        if st.button(
            "🔄 画面を更新する",
            use_container_width=True,
            key=f"btn_reload_{code}"
        ):
            st.rerun()

# ===== 追加：Pdシート生成関数 =====
def build_pd_sheet(wb, protocol_no, basic_data, pd_data, master_data):
    """
    Pd欄シートを生成してワークブックに追加する。
    O欄シートの次（index=2）に挿入。

    データフロー：
    1. 基本情報L列からPdカテゴリIDリストを取得
    2. PdシートからID照合→説明文・種別・優先順位を取得
    3. 種別A/B/C別に分類
    4. 説明文完全一致で重複統合、薬剤名を半角カナで「・」結合
    5. カテゴリ別ブロック形式で出力
    """

    # ---------- スタイル定義 ----------
    FILL_PDA_HEADER = PatternFill('solid', fgColor='4CAF50')   # 🟢 緑
    FILL_PDB_HEADER = PatternFill('solid', fgColor='FFC107')   # 🟡 黄
    FILL_PDC_HEADER = PatternFill('solid', fgColor='2196F3')   # 🔵 青
    FILL_PDA_ROW    = PatternFill('solid', fgColor='F1F8E9')   # 薄緑
    FILL_PDB_ROW    = PatternFill('solid', fgColor='FFFDE7')   # 薄黄
    FILL_PDC_ROW    = PatternFill('solid', fgColor='E3F2FD')   # 薄青
    FILL_TITLE      = PatternFill('solid', fgColor='2E4057')   # ダーク
    FILL_GUIDE      = PatternFill('solid', fgColor='FFF9C4')   # ガイド背景

    FONT_BASE   = "BIZ UDゴシック"
    FONT_WHITE  = Font(name=FONT_BASE, color='FFFFFF', bold=True, size=11)
    FONT_BOLD   = Font(name=FONT_BASE, color='000000', bold=True)
    FONT_NORMAL = Font(name=FONT_BASE, color='000000', size=10)
    FONT_GUIDE  = Font(name=FONT_BASE, color='795548', size=9)
    FONT_DRUG   = Font(name=FONT_BASE, color='1565C0', bold=True, size=10)

    thin   = Side(style='thin',   color='CCCCCC')
    medium = Side(style='medium', color='888888')
    BORDER        = Border(left=thin,   right=thin,   top=thin,   bottom=thin)
    BORDER_MEDIUM = Border(left=medium, right=medium, top=medium, bottom=medium)

    CATEGORY_CONFIG = {
        'A': {
            'label'      : '🟢 A（PDA）　毎回説明する標準説明文',
            'fill_header': FILL_PDA_HEADER,
            'fill_row'   : FILL_PDA_ROW,
        },
        'B': {
            'label'      : '🟡 B（PDB）　副作用が出た患者だけに説明',
            'fill_header': FILL_PDB_HEADER,
            'fill_row'   : FILL_PDB_ROW,
        },
        'C': {
            'label'      : '🔵 C（PDC）　特定薬剤に必ず伝える固有注意',
            'fill_header': FILL_PDC_HEADER,
            'fill_row'   : FILL_PDC_ROW,
        },
    }

    # ---------- 1. 基本情報L列からPdカテゴリIDを取得 ----------
    basic_row = next(
        (b for b in basic_data if b['プロトコールNo'] == protocol_no),
        None
    )
    if basic_row is None:
        return  # レジメンが見つからなければスキップ

    pd_cat_raw = str(basic_row.get('Pdカテゴリ', '')).strip()
    if not pd_cat_raw:
        return  # Pdカテゴリ未設定ならスキップ

    # カンマ・読点・スペース区切りに対応
    pd_ids = [x.strip() for x in re.split(r'[,、|\s]+', pd_cat_raw) if x.strip()]

    # ---------- 2. PdシートからID照合 ----------
    pd_dict = {str(p['カテゴリID']).strip(): p for p in pd_data if p.get('カテゴリID')}

    matched_pd = []
    for pid in pd_ids:
        if pid in pd_dict:
            matched_pd.append(pd_dict[pid])

    if not matched_pd:
        return  # 対応するPdデータがなければスキップ

    # ---------- 3. 薬品マスタから半角カナ辞書を作成 ----------
    # キー：一般名（全角）→ 値：一般名（半角カナ）
    master_kana = {
        str(m.get('一般名（全角）', '')).strip(): str(m.get('一般名（半角カナ）', '')).strip()
        for m in master_data
        if m.get('一般名（全角）')
    }

    def get_half_kana_name(full_name):
        """全角薬剤名→半角カナ。マスタになければto_half_kanaで変換"""
        full_name = str(full_name).strip()
        return master_kana.get(full_name) or to_half_kana(full_name)

    # ---------- 4. 種別A/B/C別に分類・優先順位順ソート ----------
    categorized = {'A': [], 'B': [], 'C': []}
    for pd_row in matched_pd:
        kubun = str(pd_row.get('種別', '')).strip().upper()
        if kubun in categorized:
            categorized[kubun].append(pd_row)

    # 優先順位（F列）でソート。空白は末尾
    for kubun in categorized:
        categorized[kubun].sort(
            key=lambda x: int(x['優先順位']) if str(x.get('優先順位','')).isdigit() else 999
        )

    # ---------- 5. 説明文重複統合（完全一致）＋薬剤名結合 ----------
    # トリガーキーワードからマスタを逆引きして薬剤名を取得する補助関数
    def get_drug_names_for_pd(pd_row):
        """
        PdシートのD列トリガーキーワード（|区切り）と
        薬品マスタのトリガーキーワードを照合して半角カナ薬剤名リストを返す。
        「手動設定」「全レジメン共通」の場合は空リストを返す。
        """
        trigger = str(pd_row.get('トリガーキーワード', '')).strip()
        if not trigger or trigger in ('手動設定', '全レジメン共通', '抗Pd'):
            return []
        keywords = [k.strip() for k in trigger.split('|') if k.strip()]
        # 薬品マスタの一般名（全角）とキーワードを照合
        matched_names = []
        for m in master_data:
            gen_name = str(m.get('一般名（全角）', '')).strip()
            if any(kw in gen_name for kw in keywords):
                kana = str(m.get('一般名（半角カナ）', '')).strip()
                if kana and kana not in matched_names:
                    matched_names.append(kana)
        return matched_names

    def merge_pd_items(pd_list):
        """
        説明文完全一致で統合。
        統合時は薬剤名を「ｼｽﾌﾟﾗﾁﾝ・ｶﾙﾎﾞﾌﾟﾗﾁﾝ：」形式に結合。
        returns: list of dict {drug_names: [...], text: str, pd_row: dict}
        """
        text_to_item = {}  # 説明文 → {drug_names, text, pd_row}
        for pd_row in pd_list:
            text = str(pd_row.get('説明文', '')).strip()
            drug_names = get_drug_names_for_pd(pd_row)
            if text in text_to_item:
                # 重複：薬剤名をマージ（重複しないように）
                for n in drug_names:
                    if n not in text_to_item[text]['drug_names']:
                        text_to_item[text]['drug_names'].append(n)
            else:
                text_to_item[text] = {
                    'drug_names': drug_names,
                    'text'      : text,
                    'pd_row'    : pd_row,
                }
        return list(text_to_item.values())

    # ---------- 6. シート生成 ----------
    ws_pd = wb.create_sheet("Pd欄")

    # O欄の次（index=2）に移動
    # wb.sheetnames = ['入力', 'O欄', 'Pd欄', '投与量シール', '説明書'] にする
    o_index = wb.sheetnames.index('O欄') if 'O欄' in wb.sheetnames else 1
    wb.move_sheet("Pd欄", offset=-(len(wb.sheetnames) - 1 - o_index))

    # 列幅設定
    ws_pd.column_dimensions['A'].width = 35   # 薬剤名列
    ws_pd.column_dimensions['B'].width = 80   # 説明文列

    # タイトル行
    ws_pd.merge_cells('A1:B1')
    ws_pd['A1'] = f'【{protocol_no}】Pd説明文　コピペ用シート'
    ws_pd['A1'].fill      = FILL_TITLE
    ws_pd['A1'].font      = Font(color='FFFFFF', bold=True, size=13)
    ws_pd['A1'].alignment = Alignment(horizontal='left', vertical='center')
    ws_pd.row_dimensions[1].height = 28

    # 使い方ガイド行
    ws_pd.merge_cells('A2:B2')
    ws_pd['A2'] = '💡 説明文をコピーして使用してください。複数薬剤で共通の説明は1行にまとめています。'
    ws_pd['A2'].fill      = FILL_GUIDE
    ws_pd['A2'].font      = FONT_GUIDE
    ws_pd['A2'].alignment = Alignment(horizontal='left', vertical='center')
    ws_pd.row_dimensions[2].height = 18

    current_row = 3

    # カテゴリA→B→C の順でブロック出力
    for kubun in ['A', 'B', 'C']:
        pd_list = categorized[kubun]
        config  = CATEGORY_CONFIG[kubun]

        # カテゴリヘッダー（空でも出力）
        current_row += 1  # 空白行
        ws_pd.merge_cells(f'A{current_row}:B{current_row}')
        ws_pd[f'A{current_row}']           = config['label']
        ws_pd[f'A{current_row}'].fill      = config['fill_header']
        ws_pd[f'A{current_row}'].font      = FONT_WHITE
        ws_pd[f'A{current_row}'].alignment = Alignment(
            horizontal='left', vertical='center'
        )
        ws_pd[f'A{current_row}'].border    = BORDER_MEDIUM
        ws_pd.row_dimensions[current_row].height = 22
        current_row += 1

        if not pd_list:
            # 該当なし行
            ws_pd.merge_cells(f'A{current_row}:B{current_row}')
            ws_pd[f'A{current_row}']           = '（このレジメンには該当する説明文がありません）'
            ws_pd[f'A{current_row}'].font      = Font(color='999999', italic=True, size=9)
            ws_pd[f'A{current_row}'].alignment = Alignment(horizontal='left', vertical='center')
            ws_pd[f'A{current_row}'].border    = BORDER
            ws_pd.row_dimensions[current_row].height = 16
            current_row += 1
            continue

        # 列ヘッダー（薬剤名 | 説明文）
        for col, label in [(1, '薬剤名'), (2, '説明文')]:
            c = ws_pd.cell(row=current_row, column=col)
            c.value     = label
            c.fill      = PatternFill('solid', fgColor='546E7A')
            c.font      = Font(color='FFFFFF', bold=True, size=10)
            c.alignment = Alignment(horizontal='center', vertical='center')
            c.border    = BORDER
        ws_pd.row_dimensions[current_row].height = 16
        current_row += 1

        # 重複統合してデータ行を出力
        merged_items = merge_pd_items(pd_list)

        for item in merged_items:
            drug_names = item['drug_names']
            text       = item['text']
            fill_row   = config['fill_row']

            # 薬剤名セル
            # 「ｼｽﾌﾟﾗﾁﾝ・ｶﾙﾎﾞﾌﾟﾗﾁﾝ：」形式
            if drug_names:
                drug_cell_val = '・'.join(drug_names) + '：'
            else:
                drug_cell_val = ''  # 手動設定・全レジメン共通は空欄

            c_drug = ws_pd.cell(row=current_row, column=1)
            c_drug.value     = drug_cell_val
            c_drug.fill      = fill_row
            c_drug.font      = FONT_DRUG
            c_drug.alignment = Alignment(
                horizontal='left', vertical='top', wrap_text=True
            )
            c_drug.border    = BORDER

            # 説明文セル（\r\n → 改行対応）
            text_clean = text.replace('\r\n', '\n').replace('\r', '\n')

            c_text = ws_pd.cell(row=current_row, column=2)
            c_text.value     = text_clean
            c_text.fill      = fill_row
            c_text.font      = FONT_NORMAL
            c_text.alignment = Alignment(
                horizontal='left', vertical='top', wrap_text=True
            )
            c_text.border    = BORDER

            # 行高さ：説明文の行数に応じて自動調整
            line_count = max(text_clean.count('\n') + 1, 1)
            ws_pd.row_dimensions[current_row].height = max(18, line_count * 15 + 5)

            current_row += 1

    return ws_pd


def create_excel(protocol_no, basic_data, drug_data,
                 master_data, notes_data,
                 pd_data=None):          # ===== 追加：pd_data引数 =====

    result = get_regimen(protocol_no, basic_data, drug_data, master_data)
    if result is None:
        return None
    basic       = result['basic']
    drugs       = result['drugs']
    master_dict = result['master_dict']

    FILL_INPUT    = PatternFill('solid', fgColor='FFFF99')
    FILL_CALC     = PatternFill('solid', fgColor='CCFFCC')
    FILL_DISABLED = PatternFill('solid', fgColor='CCCCCC')
    FILL_HEADER   = PatternFill('solid', fgColor='2E4057')
    FILL_TITLE    = PatternFill('solid', fgColor='4472C4')
    FILL_CANCER   = PatternFill('solid', fgColor='FFFDE7')
    FILL_SUPPORT  = PatternFill('solid', fgColor='D6EAF8')
    FILL_REST     = PatternFill('solid', fgColor='F2F3F4')
    FILL_RESTM    = PatternFill('solid', fgColor='EBEBEB')
    FILL_ORAL     = PatternFill('solid', fgColor='F0FFF0')
    FILL_WARNING  = PatternFill('solid', fgColor='FFF3CD')
    FONT_WHITE        = Font(name="BIZ UDゴシック", color='FFFFFF', bold=True, size=12)
    FONT_BOLD        = Font(name="BIZ UDゴシック", bold=True, size=11)
    FONT_NORMAL       = Font(name="BIZ UDゴシック", color='000000')
    FONT_DISABLED     = Font(name="BIZ UDゴシック", color='999999')
    FONT_WARNING      = Font(name="BIZ UDゴシック", color='856404', bold=True)
    FONT_WARNING_BODY = Font(name="BIZ UDゴシック", color='856404')
                     
    # ===== 追加：BIZ UDゴシック 全体フォント定義 =====
    FONT_BASE = "BIZ UDゴシック"

    # ===== 既存フォント定義を上書き =====
    FONT_WHITE        = Font(name=FONT_BASE, color='FFFFFF', bold=True)
    FONT_BOLD         = Font(name=FONT_BASE, color='000000', bold=True)
    FONT_NORMAL       = Font(name=FONT_BASE, color='000000')
    FONT_DISABLED     = Font(name=FONT_BASE, color='999999')
    FONT_WARNING      = Font(name=FONT_BASE, color='856404', bold=True)
    FONT_WARNING_BODY = Font(name=FONT_BASE, color='856404')

    thin   = Side(style='thin',   color='CCCCCC')
    thin   = Side(style='thin',   color='CCCCCC')
    medium = Side(style='medium', color='333333')
    BORDER        = Border(left=thin,   right=thin,   top=thin,   bottom=thin)
    BORDER_MEDIUM = Border(left=medium, right=medium, top=medium, bottom=medium)

    cancer_drugs = [d for d in drugs
                    if str(d.get('①O欄_抗がん剤','')) == '○']
    support_inj_all = sorted(
        [d for d in drugs
         if str(d.get('①O欄_支持療法','')) == '○'
         and str(d.get('投与順序','')) != '内服'],
        key=lambda d: INJECTION_ORDER.get(str(d.get('支持療法分類','')), 99)
    )
    support_oral_all = [d for d in drugs
                        if str(d.get('①O欄_支持療法','')) == '○'
                        and str(d.get('投与順序','')) == '内服']

    need_bsa = any(str(d['用量根拠']) == 'BSA依存' for d in cancer_drugs)
    need_ccr = any(str(d['用量根拠']) == 'AUC依存' for d in cancer_drugs)
    need_bw  = any(str(d['用量根拠']) in ('BW依存','AUC依存') for d in cancer_drugs)
    need_age = need_ccr
    need_sex = need_ccr

    wb = Workbook()
    wb.remove(wb.active)

    # ===== 入力シート（既存のまま）=====
    ws1 = wb.create_sheet("入力")
    ws1.column_dimensions['A'].width = 24
    ws1.column_dimensions['B'].width = 18
    ws1.column_dimensions['C'].width = 14
    ws1.column_dimensions['D'].width = 14
    ws1.merge_cells('A1:D1')
    ws1['A1'] = f"【{protocol_no}】{basic['レジメン名']}"
    ws1['A1'].fill      = FILL_HEADER
    ws1['A1'].font      = Font(color='FFFFFF', bold=True, size=12)
    ws1['A1'].alignment = Alignment(horizontal='left', vertical='center')
    ws1.row_dimensions[1].height = 25
    ws1.merge_cells('A2:D2')
    ws1['A2'] = f"1コース{basic['1コース日数']}日　{basic['対象疾患']}"
    ws1['A2'].fill      = FILL_TITLE
    ws1['A2'].font      = FONT_WHITE
    ws1['A2'].alignment = Alignment(horizontal='left')

    row = 4
    ws1.merge_cells(f'A{row}:D{row}')
    ws1[f'A{row}'] = '【患者情報】'
    ws1[f'A{row}'].font = FONT_BOLD
    ws1[f'A{row}'].fill = PatternFill('solid', fgColor='E8F4FD')
    row += 1

    def add_input_row(ws, row, label, needed, formula=None):
        ws[f'A{row}'] = label
        ws[f'A{row}'].font   = FONT_BOLD
        ws[f'A{row}'].border = BORDER
        if needed:
            if formula:
                ws[f'B{row}'] = formula
                ws[f'B{row}'].fill = FILL_CALC
            else:
                ws[f'B{row}'].fill = FILL_INPUT
        else:
            ws[f'B{row}'].fill = FILL_DISABLED
            ws[f'B{row}'].font = FONT_DISABLED
        ws[f'B{row}'].border    = BORDER
        ws[f'B{row}'].alignment = Alignment(horizontal='center')
        return row + 1

    row = add_input_row(ws1, row, '体重 (kg)', need_bw)
    BW_ROW = row - 1
    row = add_input_row(ws1, row, 'BSA (m2)', need_bsa)
    BSA_ROW = row - 1
    row = add_input_row(ws1, row, 'SCr', need_ccr)
    SCR_ROW = row - 1
    row = add_input_row(ws1, row, '年齢', need_age)
    AGE_ROW = row - 1

    ws1[f'A{row}'] = '性別'
    ws1[f'A{row}'].font   = FONT_BOLD
    ws1[f'A{row}'].border = BORDER
    SEX_ROW = row
    if need_sex:
        ws1[f'B{row}'].fill      = FILL_INPUT
        ws1[f'B{row}'].border    = BORDER
        ws1[f'B{row}'].alignment = Alignment(horizontal='center')
        dv = DataValidation(type='list', formula1='"男,女"',
                            allow_blank=True, showDropDown=False)
        dv.sqref = f'B{row}'
        ws1.add_data_validation(dv)
    else:
        ws1[f'B{row}'].fill   = FILL_DISABLED
        ws1[f'B{row}'].font   = FONT_DISABLED
        ws1[f'B{row}'].border = BORDER
    row += 1

    if need_ccr:
        ccr_f = (
            f'=IF(OR(B{BW_ROW}="",B{AGE_ROW}="",'
            f'B{SCR_ROW}="",B{SEX_ROW}=""),"",ROUND('
            f'IF(B{SEX_ROW}="女",'
            f'(140-B{AGE_ROW})*B{BW_ROW}/(72*B{SCR_ROW})*0.85,'
            f'(140-B{AGE_ROW})*B{BW_ROW}/(72*B{SCR_ROW})),1))'
        )
        row = add_input_row(ws1, row, 'Ccr (mL/min)', True, formula=ccr_f)
    else:
        row = add_input_row(ws1, row, 'Ccr (mL/min)', False)
    CCR_ROW = row - 1

    ws1[f'A{row}'] = '開始日'
    ws1[f'A{row}'].font          = FONT_BOLD
    ws1[f'A{row}'].border        = BORDER
    ws1[f'B{row}'].fill          = FILL_INPUT
    ws1[f'B{row}'].border        = BORDER
    ws1[f'B{row}'].number_format = 'YYYY/M/D'
    ws1[f'B{row}'].alignment     = Alignment(horizontal='center')
    START_DATE_ROW = row
    row += 1

    ws1[f'A{row}'] = 'コース目'
    ws1[f'A{row}'].font      = FONT_BOLD
    ws1[f'A{row}'].border    = BORDER
    ws1[f'B{row}'].fill      = FILL_INPUT
    ws1[f'B{row}'].border    = BORDER
    ws1[f'B{row}'].alignment = Alignment(horizontal='center')
    COURSE_ROW = row
    row += 2

    ws1.merge_cells(f'A{row}:D{row}')
    ws1[f'A{row}'] = '【投与量】'
    ws1[f'A{row}'].font = FONT_BOLD
    ws1[f'A{row}'].fill = PatternFill('solid', fgColor='E8F4FD')
    row += 1

    for col, h in enumerate(['薬剤名','計算値','投与量(mg)','達成率'], 1):
        c = ws1.cell(row=row, column=col)
        c.value = h; c.fill = FILL_HEADER; c.font = FONT_WHITE
        c.alignment = Alignment(horizontal='center'); c.border = BORDER
    row += 1

    dose_rows = {}
    for drug in cancer_drugs:
        code      = str(drug['管理コード'])
        master    = master_dict.get(code, {})
        dose_base = str(drug['用量根拠'])
        dose_num  = float(drug['投与量数値']) if str(drug['投与量数値']) != '' else 0
        name      = str(drug.get('商品名','') or master.get('採用商品名（全角）', code))
        dose_str, unit_str = format_dose_text(drug)
        ws1.cell(row=row, column=1).value  = name
        ws1.cell(row=row, column=1).font   = FONT_BOLD
        ws1.cell(row=row, column=1).border = BORDER
        if dose_base == '固定用量':
            for col in [2,3,4]:
                c = ws1.cell(row=row, column=col)
                c.fill = FILL_DISABLED; c.font = FONT_DISABLED
                c.border = BORDER; c.alignment = Alignment(horizontal='center')
            ws1.cell(row=row, column=2).value = f'{dose_str}{unit_str}'
            ws1.cell(row=row, column=3).value = f'{dose_str}{unit_str}'
            ws1.cell(row=row, column=4).value = '100%'
        elif dose_base == 'BSA依存':
            ws1.cell(row=row, column=2).value = f'=IFERROR(ROUND(B{BSA_ROW}*{dose_num},1)&"mg","")'
            ws1.cell(row=row, column=2).fill = FILL_CALC; ws1.cell(row=row, column=2).border = BORDER
            ws1.cell(row=row, column=2).alignment = Alignment(horizontal='center')
            ws1.cell(row=row, column=3).fill = FILL_INPUT; ws1.cell(row=row, column=3).border = BORDER
            ws1.cell(row=row, column=4).value = f'=IFERROR(TEXT(C{row}/(B{BSA_ROW}*{dose_num})*100,"0.0")&"%","")'
            ws1.cell(row=row, column=4).fill = FILL_CALC; ws1.cell(row=row, column=4).border = BORDER
            ws1.cell(row=row, column=4).alignment = Alignment(horizontal='center')
        elif dose_base == 'AUC依存':
            ws1.cell(row=row, column=2).value = f'=IFERROR(ROUND((B{CCR_ROW}+25)*{dose_num},0)&"mg","")'
            ws1.cell(row=row, column=2).fill = FILL_CALC; ws1.cell(row=row, column=2).border = BORDER
            ws1.cell(row=row, column=2).alignment = Alignment(horizontal='center')
            ws1.cell(row=row, column=3).fill = FILL_INPUT; ws1.cell(row=row, column=3).border = BORDER
            ws1.cell(row=row, column=4).value = f'=IFERROR(TEXT(C{row}/((B{CCR_ROW}+25)*{dose_num})*100,"0.0")&"%","")'
            ws1.cell(row=row, column=4).fill = FILL_CALC; ws1.cell(row=row, column=4).border = BORDER
            ws1.cell(row=row, column=4).alignment = Alignment(horizontal='center')
        elif dose_base == 'BW依存':
            ws1.cell(row=row, column=2).value = f'=IFERROR(ROUND(B{BW_ROW}*{dose_num},1)&"mg","")'
            ws1.cell(row=row, column=2).fill = FILL_CALC; ws1.cell(row=row, column=2).border = BORDER
            ws1.cell(row=row, column=2).alignment = Alignment(horizontal='center')
            ws1.cell(row=row, column=3).fill = FILL_INPUT; ws1.cell(row=row, column=3).border = BORDER
            ws1.cell(row=row, column=4).value = f'=IFERROR(TEXT(C{row}/(B{BW_ROW}*{dose_num})*100,"0.0")&"%","")'
            ws1.cell(row=row, column=4).fill = FILL_CALC; ws1.cell(row=row, column=4).border = BORDER
            ws1.cell(row=row, column=4).alignment = Alignment(horizontal='center')
        dose_rows[code] = row
        row += 1

    row += 1
    template_notes = sorted(
        [n for n in notes_data
         if str(n.get('プロトコールNo','')) == 'テンプレート注意事項'],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )
    if template_notes:
        ws1.merge_cells(f'A{row}:D{row}')
        ws1[f'A{row}'] = '【本テンプレートシートご使用上の注意事項】'
        ws1[f'A{row}'].font = FONT_WARNING; ws1[f'A{row}'].fill = FILL_WARNING
        ws1[f'A{row}'].alignment = Alignment(horizontal='left', vertical='center')
        ws1.row_dimensions[row].height = 18
        row += 1
        for note in template_notes:
            text = str(note['注意事項文章'])
            ws1.merge_cells(f'A{row}:D{row}')
            c = ws1.cell(row=row, column=1)
            c.value = text; c.font = FONT_WARNING_BODY; c.fill = FILL_WARNING
            c.alignment = Alignment(wrap_text=True, vertical='center', horizontal='left')
            c.border = Border(
                left=Side(style='thin', color='F0AD4E'),
                right=Side(style='thin', color='F0AD4E'),
                top=Side(style='thin', color='F0AD4E'),
                bottom=Side(style='thin', color='F0AD4E'),
            )
            ws1.row_dimensions[row].height = 40 if len(text) > 50 else 20
            row += 1

    # ===== O欄シート（既存のまま）=====
    ws2 = wb.create_sheet("O欄")
    ws2.column_dimensions['A'].width = 70
    ws2['A1'] = '▼コピーしてカルテに貼り付けてください'
    ws2['A1'].fill = PatternFill('solid', fgColor='ADD8E6')
    ws2['A1'].font = Font(bold=True)

    bsa_ref = f'入力!B{BSA_ROW}' if need_bsa else None
    bw_ref  = f'入力!B{BW_ROW}'  if need_bw  else None
    scr_ref = f'入力!B{SCR_ROW}' if need_ccr else None
    ccr_ref = f'入力!B{CCR_ROW}' if need_ccr else None
    line1   = f'"【{protocol_no}】{basic["レジメン名"]}(1ｸｰﾙ{basic["1コース日数"]}日)"'

    patient_parts = []
    if bsa_ref: patient_parts.append(f'"BSA："&TEXT({bsa_ref},"0.000")&" m2"')
    if bw_ref:  patient_parts.append(f'"BW："&TEXT({bw_ref},"0.0")&" kg"')
    if scr_ref: patient_parts.append(f'"SCr："&{scr_ref}')
    if ccr_ref: patient_parts.append(f'"Ccr："&TEXT({ccr_ref},"0.0")&" mL/min"')
    line2 = '&"  "&'.join(patient_parts) if patient_parts else None

    drug_lines = []
    for drug in cancer_drugs:
        code      = str(drug['管理コード'])
        master    = master_dict.get(code, {})
        dose_base = str(drug['用量根拠'])
        dose_num  = float(drug['投与量数値']) if str(drug['投与量数値']) != '' else 0
        day_str   = str(drug['投与Day文字'])
        name_half = to_half_kana(str(master.get('一般名（全角）','')))
        d_row     = dose_rows.get(code)
        dose_str, unit_str = format_dose_text(drug)
        if dose_base == '固定用量':
            line = f'"{name_half} ({dose_str}{unit_str})  投与量：{dose_str}{unit_str}  {day_str}"'
        elif dose_base == 'AUC依存':
            line = (f'"{name_half} (AUC{int(dose_num)})  投与量："'
                    f'&IF(入力!C{d_row}="","未入力",入力!C{d_row}&"mg"&"("&入力!D{d_row}&")")&"  {day_str}"'
                    ) if d_row else f'"{name_half} (AUC{int(dose_num)})  {day_str}"'
        else:
            line = (f'"{name_half} ({dose_num}{unit_str})  投与量："'
                    f'&IF(入力!C{d_row}="","未入力",入力!C{d_row}&"mg"&"("&入力!D{d_row}&")")&"  {day_str}"'
                    ) if d_row else f'"{name_half} ({dose_num}{unit_str})  {day_str}"'
        drug_lines.append(line)

    inj_parts     = [make_support_line(d) for d in support_inj_all]
    oral_parts    = [make_support_line(d) for d in support_oral_all]
    support_line1 = ('"支持療法：' + '､'.join(inj_parts) + '"') if inj_parts else None
    support_line2 = ('"　　　　　' + '､'.join(oral_parts) + '"') if oral_parts else None

    all_lines = [line1]
    if line2: all_lines += ['"  "', line2]
    all_lines += ['"  "'] + drug_lines
    if support_line1: all_lines += ['"  "', support_line1]
    if support_line2: all_lines.append(support_line2)

    ws2['A2'] = '=' + '&CHAR(10)&'.join(all_lines)
    ws2['A2'].alignment = Alignment(wrap_text=True, vertical='top')
    ws2['A2'].border    = BORDER_MEDIUM
    ws2.row_dimensions[2].height = 200

    # ===== 追加：Pd欄シートをO欄の次に生成 =====
    if pd_data:
        build_pd_sheet(wb, protocol_no, basic_data, pd_data, master_data)

    # ===== 投与量シール・説明書シート（既存のまま）=====
    ws3 = wb.create_sheet("投与量シール")
    ws3.column_dimensions['A'].width = 50
    ws3.column_dimensions['B'].width = 18
    short_name = shorten_regimen_name(basic['レジメン名'])
    srow = 1
    ws3.merge_cells(f'A{srow}:B{srow}')
    ws3[f'A{srow}'] = '●化学療法：'
    ws3[f'A{srow}'].font = Font(bold=True, size=11)
    srow += 1
    course_f = (
        f'=IFERROR("【{protocol_no}】{short_name}（1コース{basic["1コース日数"]}日）　"'
        f'&CHAR(10)&CHOOSE(入力!B{COURSE_ROW},"①","②","③","④","⑤","⑥","⑦","⑧","⑨","⑩")'
        f'&"コース目　"&TEXT(入力!B{START_DATE_ROW},"YYYY/M/D")&"～","")'
    )
    ws3.merge_cells(f'A{srow}:B{srow}')
    ws3[f'A{srow}'] = course_f
    ws3[f'A{srow}'].font = Font(bold=True)
    ws3[f'A{srow}'].alignment = Alignment(wrap_text=True)
    ws3.row_dimensions[srow].height = 45
    srow += 2
    ws3.merge_cells(f'A{srow}:B{srow}')
    ws3[f'A{srow}'] = '＜抗がん薬＞'
    ws3[f'A{srow}'].font = Font(bold=True)
    srow += 1

    cancer_day_nums = sorted(set(
        d for drug in cancer_drugs
        for d in parse_days_num(str(drug.get('投与Day数値','')))
    ))

    def make_date_formula(day_nums, start_row):
        if not day_nums:
            return f'=IFERROR(TEXT(入力!B{start_row},"YYYY/M/D"),"")'
        is_consec = (len(day_nums) > 1 and
                     all(day_nums[i+1]==day_nums[i]+1
                         for i in range(len(day_nums)-1)))
        if is_consec:
            offset = day_nums[-1] - day_nums[0]
            return (f'=IFERROR(TEXT(入力!B{start_row},"YYYY/M/D")'
                    f'&"〜"&TEXT(入力!B{start_row}+{offset},"M/D"),"")')
        elif len(day_nums) == 1:
            return f'=IFERROR(TEXT(入力!B{start_row},"YYYY/M/D"),"")'
        else:
            parts = [f'TEXT(入力!B{start_row},"YYYY/M/D")']
            for d in day_nums[1:]:
                parts.append(f'TEXT(入力!B{start_row}+{d-day_nums[0]},"M/D")')
            return '=IFERROR(' + '&", "&'.join(parts) + ',"")'

    ws3.merge_cells(f'A{srow}:B{srow}')
    ws3[f'A{srow}'] = make_date_formula(cancer_day_nums, START_DATE_ROW)
    ws3[f'A{srow}'].font = Font(bold=True)
    srow += 1

    for drug in cancer_drugs:
        code = str(drug['管理コード'])
        master = master_dict.get(code, {})
        name = str(drug.get('商品名','') or master.get('採用商品名（全角）', code))
        d_row = dose_rows.get(code)
        dose_base = str(drug['用量根拠'])
        dose_str, unit_str = format_dose_text(drug)
        ws3[f'A{srow}'] = name
        ws3[f'A{srow}'].font = FONT_NORMAL
        if dose_base == '固定用量':
            ws3[f'B{srow}'] = f'{dose_str}{unit_str}'
        elif d_row:
            ws3[f'B{srow}'] = f'=IFERROR(入力!C{d_row}&"mg","")'
        ws3[f'B{srow}'].alignment = Alignment(horizontal='right')
        srow += 1

    srow += 1
    ws3.merge_cells(f'A{srow}:B{srow}')
    ws3[f'A{srow}'] = '＜支持療法＞'
    ws3[f'A{srow}'].font = Font(bold=True)
    srow += 1
    if support_inj_all:
        ws3.merge_cells(f'A{srow}:B{srow}')
        ws3[f'A{srow}'] = '､'.join([make_support_line(d) for d in support_inj_all])
        ws3[f'A{srow}'].font = FONT_NORMAL
        ws3[f'A{srow}'].alignment = Alignment(wrap_text=True)
        ws3.row_dimensions[srow].height = 30
        srow += 1
    for drug in support_oral_all:
        ws3.merge_cells(f'A{srow}:B{srow}')
        ws3[f'A{srow}'] = make_support_line(drug)
        ws3[f'A{srow}'].font = FONT_NORMAL
        srow += 1

    ws4 = wb.create_sheet("説明書")
    ws4.page_setup.paperSize = 9
    ws4.page_setup.orientation = 'portrait'
    ws4.page_setup.fitToPage = True
    ws4.page_setup.fitToHeight = 1
    ws4.page_setup.fitToWidth  = 1

    COL_ORDER=1; COL_NAME=2; COL_DESC=3; COL_TIME=4; COL_DAY_START=5
    ws4.column_dimensions[get_column_letter(COL_ORDER)].width = 5
    ws4.column_dimensions[get_column_letter(COL_NAME)].width  = 22
    ws4.column_dimensions[get_column_letter(COL_DESC)].width  = 22
    ws4.column_dimensions[get_column_letter(COL_TIME)].width  = 8

    schedule_drugs = [d for d in drugs if str(d.get('④説明書','')) == '○']
    inj_schedule   = [d for d in schedule_drugs if str(d.get('投与順序','')) != '内服']
    oral_schedule  = [d for d in schedule_drugs if str(d.get('投与順序','')) == '内服']

    rp_groups_4 = defaultdict(list)
    for drug in inj_schedule:
        rp_groups_4[str(drug.get('投与順序',''))].append(drug)
    sorted_rps_4 = sorted(rp_groups_4.keys(),
                          key=lambda x: int(x) if x.isdigit() else 99)

    cycle = int(basic['1コース日数'])
    all_invest_days_4 = set()
    for rp in sorted_rps_4:
        all_invest_days_4.update(
            parse_days_num(rp_groups_4[rp][0].get('投与Day数値',''))
        )
    invest_days_4 = sorted(all_invest_days_4)
    columns_4 = []
    prev = 0
    for d in invest_days_4:
        if d > prev + 1:
            rs=prev+1; re=d-1
            columns_4.append({'type':'休薬中',
                               'label':f'{rs}日目' if rs==re else f'{rs}〜{re}日目',
                               'days':list(range(rs,re+1))})
        columns_4.append({'type':'投与','label':f'{d}日目','days':[d]})
        prev = d
    if prev < cycle:
        rs=prev+1; re=cycle
        columns_4.append({'type':'休薬',
                           'label':f'{rs}日目' if rs==re else f'{rs}〜{re}日目',
                           'days':list(range(rs,re+1))})

    n_rps_4  = len(sorted_rps_4)
    n_cols_4 = COL_DAY_START + len(columns_4) - 1
    for col_i, col in enumerate(columns_4):
        ws4.column_dimensions[
            get_column_letter(COL_DAY_START+col_i)
        ].width = 9 if col['type'] in ('投与','休薬中') else 12

    def get_rp_info_4(rp_drugs):
        rp_sorted = sorted(rp_drugs,
                           key=lambda d: INJECTION_ORDER.get(str(d.get('支持療法分類','')),99))
        names = [str(d.get('商品名','') or d.get('採用商品名（全角）','')) for d in rp_sorted]
        name_text = '＋'.join([n for n in names if n])
        classes = set(); kubuns = set()
        for d in rp_sorted:
            kubuns.add(str(d.get('薬剤区分','')))
            cls = str(d.get('支持療法分類',''))
            if cls: classes.add(cls)
        if '抗がん剤' in kubuns:
            desc = '抗がん剤です。'
        else:
            has_nausea  = bool(classes & {'NK1','5HT3'})
            has_allergy = bool(classes & {'ステロイド','抗アレルギー'})
            if has_nausea and has_allergy: desc = '吐き気やアレルギーを抑えるお薬です。'
            elif has_nausea:               desc = '吐き気を抑えるお薬です。'
            elif has_allergy:              desc = 'アレルギーを抑えるお薬です。'
            elif 'G-CSF' in classes:      desc = '白血球を増やすお薬です。'
            elif '利尿薬' in classes:     desc = '尿の量を増やすお薬です。'
            elif '解毒薬' in classes:     desc = '副作用を和らげるお薬です。'
            else:                         desc = '点滴のお薬です。'
        time_text  = str(rp_sorted[0].get('投与時間文字',''))
        kubun_main = '抗がん剤' if '抗がん剤' in kubuns else '支持療法'
        return name_text, desc, time_text, kubun_main

    erow = 1
    regimen_short = shorten_regimen_name(basic['レジメン名'])
    ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
    ws4[f'A{erow}'] = f"【{protocol_no}】{regimen_short}"
    ws4[f'A{erow}'].fill = FILL_HEADER
    ws4[f'A{erow}'].font = Font(color='FFFFFF', bold=True, size=13)
    ws4[f'A{erow}'].alignment = Alignment(horizontal='left', vertical='center')
    ws4.row_dimensions[erow].height = 25; erow += 1

    ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
    ws4[f'A{erow}'] = f'1コース{basic["1コース日数"]}日'
    ws4[f'A{erow}'].fill = FILL_TITLE; ws4[f'A{erow}'].font = FONT_WHITE
    ws4[f'A{erow}'].alignment = Alignment(horizontal='left'); erow += 2

    for title, height in [('◎治療スケジュール',18),('＜注射＞',16)]:
        ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
        ws4[f'A{erow}'] = title
        ws4[f'A{erow}'].font = Font(bold=True, size=11 if title=='◎治療スケジュール' else 10)
        ws4[f'A{erow}'].alignment = Alignment(horizontal='left')
        ws4.row_dimensions[erow].height = height; erow += 1

    for col, label in [(COL_ORDER,'順序'),(COL_NAME,'薬品名'),(COL_DESC,'説明'),(COL_TIME,'時間')]:
        c = ws4.cell(row=erow, column=col)
        c.value=label; c.fill=FILL_HEADER; c.font=FONT_WHITE
        c.alignment=Alignment(horizontal='center',vertical='center'); c.border=BORDER
    for col_i, col in enumerate(columns_4):
        col_num = COL_DAY_START+col_i
        bg = (FILL_HEADER if col['type']=='投与'
              else PatternFill('solid',fgColor='9E9E9E') if col['type']=='休薬中'
              else PatternFill('solid',fgColor='7F7F7F'))
        c = ws4.cell(row=erow, column=col_num)
        c.value=col['label']; c.fill=bg; c.font=FONT_WHITE
        c.alignment=Alignment(horizontal='center',vertical='center'); c.border=BORDER
    ws4.row_dimensions[erow].height = 16; erow += 1
    DATA_START_4 = erow

    for rp in sorted_rps_4:
        rp_drugs = rp_groups_4[rp]
        name_text, desc, time_text, kubun_main = get_rp_info_4(rp_drugs)
        row_fill  = FILL_CANCER if kubun_main=='抗がん剤' else FILL_SUPPORT
        drug_days = parse_days_num(rp_drugs[0].get('投与Day数値',''))
        ws4.row_dimensions[erow].height = max(18, (name_text.count('＋')+1)*14)
        for col, val, fill, align_h, wrap in [
            (COL_ORDER,rp,PatternFill('solid',fgColor='F5F5F5'),'center',False),
            (COL_NAME,name_text,row_fill,'left',True),
            (COL_DESC,desc,row_fill,'left',True),
            (COL_TIME,time_text,row_fill,'center',False),
        ]:
            c = ws4.cell(row=erow, column=col)
            c.value=val; c.fill=fill
            c.font=FONT_BOLD if col==COL_ORDER else FONT_NORMAL
            c.alignment=Alignment(horizontal=align_h,vertical='center',wrap_text=wrap)
            c.border=BORDER
        for col_i, col in enumerate(columns_4):
            col_num = COL_DAY_START+col_i
            c = ws4.cell(row=erow, column=col_num)
            if col['type'] == '投与':
                hit = any(d in drug_days for d in col['days'])
                if hit:
                    c.value='●'; c.fill=row_fill
                else:
                    c.value='－'; c.fill=PatternFill('solid',fgColor='FAFAFA')
                    c.font=Font(color='AAAAAA')
            c.alignment=Alignment(horizontal='center',vertical='center'); c.border=BORDER
        erow += 1

    for col_i, col in enumerate(columns_4):
        if col['type'] in ('休薬','休薬中'):
            col_num = COL_DAY_START+col_i
            col_letter = get_column_letter(col_num)
            start_r = DATA_START_4; end_r = DATA_START_4+n_rps_4-1
            for r in range(start_r, end_r+1):
                ws4.cell(row=r, column=col_num).value = None
            ws4.merge_cells(f'{col_letter}{start_r}:{col_letter}{end_r}')
            c = ws4.cell(row=start_r, column=col_num)
            c.value = 'お休み' if col['type']=='休薬' else '－'
            c.fill  = FILL_REST if col['type']=='休薬' else FILL_RESTM
            c.font  = Font(bold=True, color='555555')
            c.alignment = Alignment(horizontal='center',vertical='center')
            c.border = BORDER
    erow += 1

    if oral_schedule:
        ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
        ws4[f'A{erow}'] = '＜内服＞'
        ws4[f'A{erow}'].font = Font(bold=True)
        ws4[f'A{erow}'].alignment = Alignment(horizontal='left')
        ws4.row_dimensions[erow].height = 16; erow += 1
        for drug in oral_schedule:
            name   = str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
            desc   = str(drug.get('患者向け説明',''))
            timing = str(drug.get('投与タイミング',''))
            ws4.row_dimensions[erow].height = 30
            ws4.merge_cells(f'A{erow}:B{erow}')
            c = ws4.cell(row=erow, column=1)
            c.value=f'・{name}：'; c.font=FONT_BOLD; c.fill=FILL_ORAL
            c.alignment=Alignment(vertical='center',horizontal='left'); c.border=BORDER
            ws4.merge_cells(f'C{erow}:{get_column_letter(n_cols_4)}{erow}')
            c = ws4.cell(row=erow, column=3)
            oral_text = desc
            if timing and timing not in ('','ー','nan'):
                if oral_text: oral_text += '\n'
                oral_text += timing
            c.value=oral_text; c.font=FONT_NORMAL; c.fill=FILL_ORAL
            c.alignment=Alignment(wrap_text=True,vertical='center',horizontal='left')
            c.border=BORDER; erow += 1

    erow += 1
    ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
    ws4[f'A{erow}'] = '◎注意事項'
    ws4[f'A{erow}'].font = Font(bold=True, size=11)
    ws4[f'A{erow}'].alignment = Alignment(horizontal='left')
    ws4.row_dimensions[erow].height = 18; erow += 1

    common_notes = sorted(
        [n for n in notes_data if str(n.get('プロトコールNo',''))=='共通'],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )
    for note in common_notes:
        ws4.merge_cells(f'A{erow}:{get_column_letter(n_cols_4)}{erow}')
        c = ws4.cell(row=erow, column=1)
        c.value = '・' + str(note['注意事項文章'])
        c.font  = FONT_NORMAL
        c.alignment = Alignment(wrap_text=True,vertical='center',horizontal='left')
        ws4.row_dimensions[erow].height = 20; erow += 1

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output.getvalue()


# ===== create_pptx は変更なし（省略せず全文そのまま）=====
def create_pptx(protocol_no, basic_data, drug_data,
                master_data, notes_data):
    result = get_regimen(protocol_no, basic_data, drug_data, master_data)
    if result is None:
        return None
    basic       = result['basic']
    drugs       = result['drugs']
    master_dict = result['master_dict']
    cycle       = int(basic['1コース日数'])

    COLOR_HEADER  = RGBColor(0x40,0x40,0x40)
    COLOR_ICI     = RGBColor(0x70,0xAD,0x47)
    COLOR_MOL     = RGBColor(0xFF,0xE0,0x99)
    COLOR_NAUSEA  = RGBColor(0x84,0xAC,0xD4)
    COLOR_CYTO    = RGBColor(0xF4,0xB1,0x83)
    COLOR_AUX     = RGBColor(0x9D,0xD7,0xEA)
    COLOR_REST_BG = RGBColor(0xBF,0xBF,0xBF)
    COLOR_WHITE   = RGBColor(0xFF,0xFF,0xFF)
    COLOR_BLACK   = RGBColor(0x00,0x00,0x00)
    TYPE_COLOR_MAP = {
        '免疫チェックポイント阻害薬': COLOR_ICI,
        '分子標的薬': COLOR_MOL,
        '吐き気止め': COLOR_NAUSEA,
        '細胞障害性抗がん薬': COLOR_CYTO,
        '補助薬': COLOR_AUX,
        '輸液': COLOR_AUX,
        'その他支持療法': COLOR_NAUSEA,
        '造血因子': COLOR_AUX,
    }
    A4_W = Cm(25.4); A4_H = Cm(19.05)

    def rgb_hex(rgb):
        return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

    def set_cell_bg(cell, rgb):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', rgb_hex(rgb))

    def set_cell_text_multi(cell, lines, font_size=Pt(11),
                             bold=False, color=COLOR_BLACK,
                             align=PP_ALIGN.CENTER):
        tf = cell.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = str(line); run.font.size = font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold = bold; run.font.color.rgb = color

    def set_cell_margin(cell, top=Cm(0.05), bottom=Cm(0.05),
                         left=Cm(0.1), right=Cm(0.1)):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        tcPr.set('marT',str(int(top))); tcPr.set('marB',str(int(bottom)))
        tcPr.set('marL',str(int(left))); tcPr.set('marR',str(int(right)))

    def merge_cells_vertical(table, col, start_row, end_row):
        for r in range(start_row, end_row+1):
            tc = table.cell(r,col)._tc; tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:vMerge')): tcPr.remove(old)
            vMerge = etree.SubElement(tcPr, qn('a:vMerge'))
            if r == start_row: vMerge.set('val','restart')

    def add_textbox(slide, x, y, w, h, text,
                    font_size=Pt(11), bold=False,
                    color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text=str(text); run.font.size=font_size
        run.font.name = "BIZ UDゴシック"
        run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def add_textbox_multi(slide, x, y, w, h, lines,
                          font_size=Pt(11), bold=False,
                          color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text=str(line); run.font.size=font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def safe_float(val):
        try:
            return float(val) if str(val).strip() != '' else 0
        except:
            return 0

    seal_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='外来手帳シール'
         and str(n.get('プロトコールNo',''))=='共通'],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )
    tegaki_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='手帳シール但し書き'
         and str(n.get('プロトコールNo','')) in ('共通',protocol_no)],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )

    schedule_drugs = [d for d in drugs if str(d.get('④説明書',''))=='○']
    inj_schedule   = [d for d in schedule_drugs if str(d.get('投与順序',''))!='内服']
    oral_schedule  = [d for d in schedule_drugs if str(d.get('投与順序',''))=='内服']

    rp_groups = defaultdict(list)
    for drug in inj_schedule:
        rp_groups[str(drug.get('投与順序',''))].append(drug)
    sorted_rps = sorted(rp_groups.keys(),
                        key=lambda x: int(x) if x.isdigit() else 99)

    all_invest_days = set()
    for rp in sorted_rps:
        all_invest_days.update(
            parse_days_num(rp_groups[rp][0].get('投与Day数値',''))
        )
    invest_days = sorted(all_invest_days)
    columns = []; prev = 0
    for d in invest_days:
        if d > prev+1:
            rs=prev+1; re=d-1
            columns.append({'type':'休薬中',
                             'label':f'{rs}日目' if rs==re else f'{rs}〜{re}日目',
                             'days':list(range(rs,re+1))})
        columns.append({'type':'投与','label':f'{d}日目','days':[d]})
        prev = d
    if prev < cycle:
        rs=prev+1; re=cycle
        columns.append({'type':'休薬','label':f'〜{re}日目',
                        'days':list(range(rs,re+1))})

    need_bw  = any(str(d.get('用量根拠','')) in ('BW依存','AUC依存')
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')
    need_bsa = any(str(d.get('用量根拠',''))=='BSA依存'
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')

    rp_all_groups = defaultdict(list)
    for drug in drugs:
        rp = str(drug.get('投与順序',''))
        if rp != '内服':
            rp_all_groups[rp].append(drug)

    total_min = int(sum(
        max((safe_float(d.get('投与時間数値', 0))
             for d in rp_drugs), default=0)
        for rp_drugs in rp_all_groups.values()
    ) * 60)

    def get_rp_info(rp_drugs):
        rp_sorted = sorted(rp_drugs,
                           key=lambda d: INJECTION_ORDER.get(str(d.get('支持療法分類','')),99))
        names = [str(d.get('商品名','') or d.get('採用商品名（全角）','')) for d in rp_sorted]
        name_text = '\n'.join([n for n in names if n])
        dose_parts = []
        for d in rp_sorted:
            dose_base = str(d.get('用量根拠',''))
            kubun     = str(d.get('薬剤区分',''))
            try:    dose_num = float(d.get('投与量数値',0) or 0)
            except: dose_num = 0
            dose_str, unit_str = format_dose_text(d)
            if kubun == '抗がん剤':
                if dose_base=='BSA依存':   dose_parts.append(f'({dose_num}mg/m²)')
                elif dose_base=='AUC依存': dose_parts.append(f'(AUC：{int(dose_num)})')
                elif dose_base=='BW依存':  dose_parts.append(f'({dose_num}mg/kg)')
                elif dose_base=='固定用量':dose_parts.append(f'({dose_str}{unit_str})')
        time_val  = rp_sorted[0].get('投与時間文字','')
        time_text = f'(投与時間：{time_val})' if time_val else ''
        type_str  = ''
        for d in rp_sorted:
            t = str(master_dict.get(str(d.get('管理コード','')),{}).get('スケジュールシール用種類',''))
            if t: type_str = t; break
        drug_days = parse_days_num(rp_sorted[0].get('投与Day数値',''))
        return name_text, dose_parts, time_text, type_str, drug_days

    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mg=Cm(0.8); cw=A4_W-mg*2; x0=mg; y0=mg
    title_w=cw*0.72; date_w=cw*0.28
    title_h=Cm(1.5); course_h=Cm(1.0)

    add_textbox(slide,x0,y0,title_w,title_h,
                text=basic['レジメン名'],
                font_size=Pt(24),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0,y0+title_h,title_w,course_h,
                text=f'（1コース{cycle}日）',
                font_size=Pt(20),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0+title_w,y0,date_w,title_h+course_h,
                text='当レジメンの開始日\n年　　月　　日',
                font_size=Pt(11),bold=False,color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    y_cur = y0+title_h+course_h+Cm(0.2)
    note_lines = (['・'+str(n['注意事項文章']) for n in seal_notes] if seal_notes else [
        '・薬剤の他、副作用を予防する、投与ラインを満たす、薬剤を洗い流すために生理食塩液等を投与します。',
        '・強い副作用が出た場合や患者様の状態によっては、治療を延期したり、減量することがあります。'
    ])
    note_h = Cm(0.55*len(note_lines))
    add_textbox_multi(slide,x0,y_cur,cw,note_h,
                      lines=note_lines,font_size=Pt(11),
                      color=COLOR_BLACK,align=PP_ALIGN.LEFT)
    y_cur += note_h+Cm(0.2)

    has_right   = bool(oral_schedule) or need_bw or need_bsa
    tbl_w_ratio = 0.60 if has_right else 1.0
    right_w     = cw*0.38
    right_x     = x0+cw*tbl_w_ratio+cw*0.02

    type_w_cm=2.2; name_w_cm=3.8; inv_w_cm=1.1; rest_w_cm=1.5
    col_w_cm = [inv_w_cm if c['type']=='投与'
                else inv_w_cm*0.8 if c['type']=='休薬中'
                else rest_w_cm for c in columns]
    tbl_total_w  = type_w_cm+name_w_cm+sum(col_w_cm)
    scale        = float((cw*tbl_w_ratio)/Cm(tbl_total_w))
    type_w_cm   *= scale; name_w_cm *= scale
    col_w_cm     = [w*scale for w in col_w_cm]

    n_cols=2+len(columns); n_rps=len(sorted_rps); n_rows=1+n_rps

    def calc_lines(nd,dp,tt):
        return nd.count('\n')+1+len(dp)+(1 if tt else 0)

    row_heights = [Cm(0.65)]
    for rp in sorted_rps:
        nd,dp,tt,_,_ = get_rp_info(rp_groups[rp])
        row_heights.append(max(Cm(0.75),Cm(0.42*calc_lines(nd,dp,tt)+0.2)))
    tbl_h = sum(row_heights)

    shape = slide.shapes.add_table(n_rows,n_cols,x0,y_cur,
                                   Cm(tbl_total_w*scale),tbl_h)
    table = shape.table
    table.columns[0].width = Cm(type_w_cm)
    table.columns[1].width = Cm(name_w_cm)
    for i,w in enumerate(col_w_cm): table.columns[2+i].width = Cm(w)
    for i,h in enumerate(row_heights): table.rows[i].height = h

    for col_i,label in enumerate(['種類','薬剤名']):
        cell = table.cell(0,col_i)
        set_cell_text_multi(cell,[label],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)
    for i,col in enumerate(columns):
        cell = table.cell(0,2+i)
        set_cell_text_multi(cell,[col['label']],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)

    rest_col_indices = [(i,col) for i,col in enumerate(columns)
                        if col['type'] in ('休薬','休薬中')]
    for row_i,rp in enumerate(sorted_rps):
        r=row_i+1
        name_text,dose_parts,time_text,type_str,drug_days = get_rp_info(rp_groups[rp])
        row_color = TYPE_COLOR_MAP.get(type_str,COLOR_NAUSEA)
        cell = table.cell(r,0)
        set_cell_text_multi(cell,[type_str],font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        name_lines = name_text.split('\n')+dose_parts
        if time_text: name_lines.append(time_text)
        cell = table.cell(r,1)
        set_cell_text_multi(cell,name_lines,font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        for i,col in enumerate(columns):
            cell = table.cell(r,2+i); set_cell_margin(cell)
            if col['type']=='投与':
                hit = any(d in drug_days for d in col['days'])
                if hit:
                    set_cell_text_multi(cell,['●'],font_size=Pt(14),bold=True,
                                         color=COLOR_BLACK,align=PP_ALIGN.CENTER)
                    set_cell_bg(cell,row_color)
                else:
                    set_cell_bg(cell,COLOR_REST_BG)

    for i,col in rest_col_indices:
        col_idx = 2+i
        merge_cells_vertical(table,col_idx,1,n_rps)
        cell = table.cell(1,col_idx)
        set_cell_bg(cell,COLOR_REST_BG)
        set_cell_text_multi(cell,[''],font_size=Pt(11),
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_margin(cell)

    tbl_bottom  = y_cur+tbl_h
    y_after_tbl = tbl_bottom+Cm(0.15)
    if tegaki_notes:
        nb_h = Cm(0.5*len(tegaki_notes))
        add_textbox_multi(slide,x0,y_after_tbl,cw*tbl_w_ratio,nb_h,
                          lines=[str(n['注意事項文章']) for n in tegaki_notes],
                          font_size=Pt(11),color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    if has_right:
        ry = y_cur
        if oral_schedule:
            oral_lines = []
            for drug in oral_schedule:
                name   = str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
                timing = str(drug.get('投与タイミング',''))
                line   = f'・{name}'
                if timing and timing not in ('','ー','nan'):
                    line += f'\n　{timing}'
                oral_lines.append(line)
            oral_box_h = Cm(0.55*sum(l.count('\n')+1 for l in oral_lines)+0.3)
            add_textbox_multi(slide,right_x+Cm(0.15),ry+Cm(0.1),
                              right_w-Cm(0.3),oral_box_h,
                              lines=oral_lines,font_size=Pt(11),
                              color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += oral_box_h+Cm(0.3)
        if need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(0.7),
                        text='体重　　　　　kg',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(0.9)
        if need_bsa or need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(1.0),
                        text='体表面積\n\n＿＿＿＿m²',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(1.3)
        add_textbox(slide,right_x,ry,right_w,Cm(0.8),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        ry += Cm(1.0)
        add_textbox(slide,right_x,ry,right_w,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
    else:
        footer_y = A4_H-Cm(1.5)-mg
        add_textbox(slide,x0,footer_y,cw*0.5,Cm(0.7),
                    text='体表面積　　　　m²',font_size=Pt(11),
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        add_textbox(slide,x0+cw*0.5,footer_y,cw*0.5,Cm(0.7),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
        add_textbox(slide,x0,footer_y+Cm(0.8),cw,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ===== Streamlit UI =====
st.title("💊 テンプレート生成")
st.caption("登録済みレジメンから Excel・パワポを 生成")
st.divider()

with st.spinner("スプレッドシートからデータを読み込み中..."):
    basic_data, drug_data, master_data, notes_data, pd_data, ae_data = load_all_data()

if not basic_data:
    st.error("データの読み込みに失敗しました。スプレッドシートの公開設定を確認してください。")
    st.stop()

regimen_list = [
    f"{b['プロトコールNo']}　{b['レジメン名']}"
    for b in basic_data
    if b.get('プロトコールNo','').strip()
]

st.subheader("📋 レジメン選択")

if not regimen_list:
    st.error("レジメン一覧が取得できませんでした。")
    st.stop()

selected = st.selectbox(
    "生成するレジメンを選択してください",
    options=regimen_list,
    index=0
)

if selected is None:
    st.stop()

protocol_no = selected.split('　')[0].strip()

selected_basic = next(
    (b for b in basic_data if b['プロトコールNo'] == protocol_no),
    None
)
if selected_basic:
    col1, col2, col3 = st.columns(3)
    col1.metric("プロトコールNo", protocol_no)
    col2.metric("1コース日数", f"{selected_basic.get('1コース日数','?')}日")
    col3.metric("対象疾患", selected_basic.get('対象疾患','?'))

    pd_cats = str(selected_basic.get('Pdカテゴリ', '')).strip()
    if not pd_cats:
        st.warning("⚠️ このレジメンのPdカテゴリが未設定です。")
        col_pd1, col_pd2 = st.columns(2)
        with col_pd1:
            if st.button(
                "📝 Pd説明文管理ページで設定する",
                use_container_width=True
            ):
                st.switch_page("pages/4_Pd説明文管理.py")
        with col_pd2:
            if st.button(
                "⏭️ スキップしてファイル生成へ",
                use_container_width=True
            ):
                st.session_state["pd_skip"] = True
                st.rerun()
    else:
        st.success(f"✅ Pdカテゴリ設定済み：{pd_cats}")

    # ===== 抗がん剤副作用マスタ登録状況チェック =====
    # このレジメンの薬剤を取得
    regimen_drugs = [
        d for d in drug_data
        if str(d.get('プロトコールNo', '')).strip() == protocol_no
    ]
    # 抗がん剤のみ対象（管理コードがACで始まるもの）
    cancer_drug_codes = list(dict.fromkeys([
        str(d.get('管理コード', '')).strip()
        for d in regimen_drugs
        if str(d.get('管理コード', '')).strip().startswith('AC')
    ]))

    if cancer_drug_codes:
        # 副作用マスタを辞書化
        ae_dict = {
            str(r.get('管理コード', '')).strip(): r
            for r in ae_data
        }
        # 登録済み・未登録に分類
        registered   = []
        unregistered = []
        for code in cancer_drug_codes:
            ae_row   = ae_dict.get(code, {})
            reg_date = str(ae_row.get('登録日', '')).strip()
            name     = str(ae_row.get('一般名（全角）', code)).strip()
            # 薬品マスタから名前を補完
            if name == code:
                name = next(
                    (m.get('一般名（全角）', code)
                     for m in master_data
                     if str(m.get('管理コード', '')).strip() == code),
                    code
                )
            if reg_date:
                registered.append({'code': code, 'name': name, 'date': reg_date})
            else:
                unregistered.append({'code': code, 'name': name})

        st.divider()
        st.subheader("📊 抗がん剤副作用マスタ 登録状況")

        if unregistered:
            st.warning(f"⚠️ 副作用が未登録の薬剤が {len(unregistered)} 件あります")

        col_reg, col_unreg = st.columns(2)

        with col_reg:
            if registered:
                st.markdown("**✅ 登録済み**")
                for r in registered:
                    st.caption(f"　{r['name']}：{r['date']} 登録")

        with col_unreg:
            if unregistered:
                st.markdown("**❌ 未登録**")
                for r in unregistered:
                    st.caption(f"　{r['name']}")

        if unregistered:
            st.divider()
            # 未登録薬剤をその場で登録するか確認
            if not st.session_state.get("ae_skip", False):
                col_go, col_skip = st.columns(2)
                with col_go:
                    if st.button(
                        "📝 このページで副作用を登録する",
                        type="primary",
                        use_container_width=True
                    ):
                        st.session_state["ae_reg_index"]      = 0
                        st.session_state["ae_reg_start"]      = True
                        st.session_state["current_protocol_no"] = protocol_no
                with col_skip:
                    if st.button(
                        "⏭️ スキップしてファイル生成へ（データ最新化）",
                        use_container_width=True,
                    ):
                        st.cache_data.clear()
                        st.session_state["ae_skip"] = True
                        st.rerun()

        # 副作用登録UI表示
        if st.session_state.get("ae_reg_start", False):
            show_ae_register_ui(
                unregistered, ae_data, master_data,
                drug_data, basic_data, pd_data
            )

st.divider()
st.subheader("📁 ファイル生成")

col_excel, col_pptx = st.columns(2)

with col_excel:
    if st.button("📊 Excel生成", use_container_width=True, type="primary"):
        with st.spinner("Excel生成中..."):
            excel_data = create_excel(
                protocol_no, basic_data,
                drug_data, master_data, notes_data,
                pd_data=pd_data
            )
        if excel_data:
            st.download_button(
                label="⬇️ Excelをダウンロード",
                data=excel_data,
                file_name=f"{protocol_no}_AI作成要チェック_{today_str}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
            st.success("✅ Excel生成完了！")
        else:
            st.error("❌ 生成に失敗しました")

with col_pptx:
    if st.button("📑 パワポ生成", use_container_width=True, type="primary"):
        with st.spinner("パワポ生成中..."):
            pptx_data = create_pptx(
                protocol_no, basic_data,
                drug_data, master_data, notes_data
            )
        if pptx_data:
            st.download_button(
                label="⬇️ パワポをダウンロード",
                data=pptx_data,
                file_name=f"{protocol_no}_スケジュールシール_AI作成要チェック_{today_str}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
            st.success("✅ パワポ生成完了！")
        else:
            st.error("❌ 生成に失敗しました")

st.divider()
st.caption("⚠️ 生成されたファイルは必ず内容を確認してから使用してください")
