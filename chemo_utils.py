"""
chemo_utils.py
化学療法テンプレートアプリ 共通ユーティリティ
"""
import re
import io
import math
from collections import defaultdict
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

INJECTION_ORDER = {
    'NK1':1,'5HT3':2,'ステロイド':3,'G-CSF':4,
    '利尿薬':5,'解毒薬':6,'抗アレルギー':7,
    'H2ブロッカー':8,'電解質補正':9,'その他注射':10
}

def to_half_kana(text):
    table = {
        'ア':'ｱ','イ':'ｲ','ウ':'ｳ','エ':'ｴ','オ':'ｵ',
        'カ':'ｶ','キ':'ｷ','ク':'ｸ','ケ':'ｹ','コ':'ｺ',
        'サ':'ｻ','シ':'ｼ','ス':'ｽ','セ':'ｾ','ソ':'ｿ',
        'タ':'ﾀ','チ':'ﾁ','ツ':'ﾂ','テ':'ﾃ','ト':'ﾄ',
        'ナ':'ﾅ','ニ':'ﾆ','ヌ':'ﾇ','ネ':'ﾈ','ノ':'ﾉ',
        'ハ':'ﾊ','ヒ':'ﾋ','フ':'ﾌ','ヘ':'ﾍ','ホ':'ﾎ',
        'マ':'ﾏ','ミ':'ﾐ','ム':'ﾑ','メ':'ﾒ','モ':'ﾓ',
        'ヤ':'ﾔ','ユ':'ﾕ','ヨ':'ﾖ',
        'ラ':'ﾗ','リ':'ﾘ','ル':'ﾙ','レ':'ﾚ','ロ':'ﾛ',
        'ワ':'ﾜ','ヲ':'ｦ','ン':'ﾝ',
        'ァ':'ｧ','ィ':'ｨ','ゥ':'ｩ','ェ':'ｪ','ォ':'ｫ',
        'ッ':'ｯ','ャ':'ｬ','ュ':'ｭ','ョ':'ｮ',
        'ガ':'ｶﾞ','ギ':'ｷﾞ','グ':'ｸﾞ','ゲ':'ｹﾞ','ゴ':'ｺﾞ',
        'ザ':'ｻﾞ','ジ':'ｼﾞ','ズ':'ｽﾞ','ゼ':'ｾﾞ','ゾ':'ｿﾞ',
        'ダ':'ﾀﾞ','ヂ':'ﾁﾞ','ヅ':'ﾂﾞ','デ':'ﾃﾞ','ド':'ﾄﾞ',
        'バ':'ﾊﾞ','ビ':'ﾋﾞ','ブ':'ﾌﾞ','ベ':'ﾍﾞ','ボ':'ﾎﾞ',
        'パ':'ﾊﾟ','ピ':'ﾋﾟ','プ':'ﾌﾟ','ペ':'ﾍﾟ','ポ':'ﾎﾟ',
        'ー':'ｰ','ヴ':'ｳﾞ','・':'･',
    }
    result = ''
    for char in str(text):
        result += table.get(char, char)
    return result



def shorten_regimen_name(regimen_name):
    name = regimen_name
    name = re.sub(r'[（(][^）)]*[）)]', '', name)
    for word in [
        '術前','術後','周術期','切除不能','再発','進行',
        '維持','補助','一次治療','二次治療','初回','難治性',
        '肺癌','胃癌','大腸癌','乳癌','膵癌','肝癌',
        '食道癌','子宮癌','卵巣癌','前立腺癌','膀胱癌',
        '腎癌','甲状腺癌','悪性リンパ腫','白血病',
        '骨髄腫','中皮腫','胸腺癌','胸腺腫','神経内分泌腫瘍',
    ]:
        name = name.replace(word, '')
    name = re.sub(r'\s+', ' ', name).strip('　 ')
    return name



def parse_days_num(day_str):
    days = []
    for part in str(day_str).split('|'):
        part = part.strip()
        if '-' in part:
            try:
                s, e = part.split('-')
                days.extend(range(int(s), int(e)+1))
            except:
                pass
        elif part.isdigit():
            days.append(int(part))
    return sorted(set(days))



def format_dose_text(drug):
    try:
        _rv = str(drug.get('投与量数値', '') or '').strip()
        _rv = ''.join(c for c in _rv if c.isdigit() or c == '.')
        dose = float(_rv or 0)
    except:
        dose = 0
    unit_input = str(drug.get('投与単位', '')).strip()
    dose_base  = str(drug.get('用量根拠', ''))
    v_to_mg    = drug.get('1V当たりmg', '')
    if unit_input.upper() == 'V':
        if v_to_mg != '' and str(v_to_mg).strip() != '':
            try:
                dose = dose * float(v_to_mg)
                unit_input = 'mg'
            except:
                pass
    if unit_input in ('mg','mg/body','mg/ body'):
        display_unit = 'mg'
    elif unit_input == '':
        display_unit = '' if dose_base == 'AUC依存' else 'mg'
    else:
        display_unit = unit_input
    dose_str = str(int(dose)) if dose == int(dose) else str(dose)
    return dose_str, display_unit

INJECTION_ORDER = {
    'NK1':1,'5HT3':2,'ステロイド':3,'G-CSF':4,
    '利尿薬':5,'解毒薬':6,'抗アレルギー':7,
    'H2ブロッカー':8,'電解質補正':9,'その他注射':10
}



def make_support_line(drug):
    name = to_half_kana(
        str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
    )
    dose_str, unit_str = format_dose_text(drug)
    day = str(drug.get('投与Day文字',''))
    return f"{name} {dose_str}{unit_str}({day})"



def get_regimen(protocol_no, basic_data, drug_data, master_data):
    master_dict = {m['管理コード']: m for m in master_data}
    basic = [b for b in basic_data
             if b['プロトコールNo'] == protocol_no]
    if not basic:
        return None
    basic = basic[0]
    drugs_raw = sorted(
        [d for d in drug_data
         if d['プロトコールNo'] == protocol_no],
        key=lambda x: (
            0 if str(x['投与順序']).isdigit() else 1,
            int(x['投与順序']) if str(x['投与順序']).isdigit() else 99
        )
    )
    drugs = []
    for drug in drugs_raw:
        code   = str(drug['管理コード'])
        master = master_dict.get(code, {})
        merged = dict(drug)
        merged.update({
            '一般名（全角）'        : master.get('一般名（全角）',''),
            '一般名（半角カナ）'    : master.get('一般名（半角カナ）',''),
            '採用商品名（全角）'    : master.get('採用商品名（全角）',''),
            '採用商品名（半角カナ）': master.get('採用商品名（半角カナ）',''),
            '薬効分類'             : master.get('薬効分類',''),
            '薬剤区分'             : master.get('薬剤区分',''),
            '支持療法分類'         : master.get('支持療法分類',''),
            '薬品マスタ単位'       : master.get('単位',''),
            '投与経路'             : master.get('投与経路',''),
            '1V当たりmg'           : master.get('1V当たりmg',''),
            '患者向け説明'         : master.get('患者向け説明',''),
            '短縮注記'             : master.get('短縮注記',''),
        })
        drugs.append(merged)
    return {'basic': basic, 'drugs': drugs, 'master_dict': master_dict}



# ===== 追加：副作用登録UI関数 =====



def create_pptx(protocol_no, basic_data, drug_data,
                master_data, notes_data):
    result = get_regimen(protocol_no, basic_data, drug_data, master_data)
    if result is None:
        return None
    basic       = result['basic']
    drugs       = result['drugs']
    master_dict = result['master_dict']
    cycle       = int(basic['1コース日数'])

    COLOR_HEADER  = RGBColor(0x40,0x40,0x40)
    COLOR_ICI     = RGBColor(0x70,0xAD,0x47)
    COLOR_MOL     = RGBColor(0xFF,0xE0,0x99)
    COLOR_NAUSEA  = RGBColor(0x84,0xAC,0xD4)
    COLOR_CYTO    = RGBColor(0xF4,0xB1,0x83)
    COLOR_AUX     = RGBColor(0x9D,0xD7,0xEA)
    COLOR_REST_BG = RGBColor(0xBF,0xBF,0xBF)
    COLOR_WHITE   = RGBColor(0xFF,0xFF,0xFF)
    COLOR_BLACK   = RGBColor(0x00,0x00,0x00)
    TYPE_COLOR_MAP = {
        '免疫チェックポイント阻害薬': COLOR_ICI,
        '分子標的薬': COLOR_MOL,
        '吐き気止め': COLOR_NAUSEA,
        '細胞障害性抗がん薬': COLOR_CYTO,
        '補助薬': COLOR_AUX,
        '輸液': COLOR_AUX,
        'その他支持療法': COLOR_NAUSEA,
        '造血因子': COLOR_AUX,
    }
    A4_W = Cm(25.4); A4_H = Cm(19.05)

    def rgb_hex(rgb):
        return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

    def set_cell_bg(cell, rgb):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', rgb_hex(rgb))

    def set_cell_text_multi(cell, lines, font_size=Pt(11),
                             bold=False, color=COLOR_BLACK,
                             align=PP_ALIGN.CENTER):
        tf = cell.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = str(line); run.font.size = font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold = bold; run.font.color.rgb = color

    def set_cell_margin(cell, top=Cm(0.05), bottom=Cm(0.05),
                         left=Cm(0.1), right=Cm(0.1)):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        tcPr.set('marT',str(int(top))); tcPr.set('marB',str(int(bottom)))
        tcPr.set('marL',str(int(left))); tcPr.set('marR',str(int(right)))

    def merge_cells_vertical(table, col, start_row, end_row):
        for r in range(start_row, end_row+1):
            tc = table.cell(r,col)._tc; tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:vMerge')): tcPr.remove(old)
            vMerge = etree.SubElement(tcPr, qn('a:vMerge'))
            if r == start_row: vMerge.set('val','restart')

    def add_textbox(slide, x, y, w, h, text,
                    font_size=Pt(11), bold=False,
                    color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text=str(text); run.font.size=font_size
        run.font.name = "BIZ UDゴシック"
        run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def add_textbox_multi(slide, x, y, w, h, lines,
                          font_size=Pt(11), bold=False,
                          color=COLOR_BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x,y,w,h)
        tf = txBox.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text=str(line); run.font.size=font_size
            run.font.name = "BIZ UDゴシック"
            run.font.bold=bold; run.font.color.rgb=color
        return txBox

    def safe_float(val):
        try:
            return float(val) if str(val).strip() != '' else 0
        except:
            return 0

    seal_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='外来手帳シール'
         and str(n.get('プロトコールNo',''))=='共通'],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )
    tegaki_notes = sorted(
        [n for n in notes_data
         if str(n.get('区分',''))=='手帳シール但し書き'
         and str(n.get('プロトコールNo','')) in ('共通',protocol_no)],
        key=lambda x: int(x['順序']) if str(x['順序']).isdigit() else 99
    )

    schedule_drugs = [d for d in drugs if str(d.get('④説明書',''))=='○']
    inj_schedule   = [d for d in schedule_drugs if str(d.get('投与順序',''))!='内服']
    oral_schedule  = [d for d in schedule_drugs if str(d.get('投与順序',''))=='内服']

    rp_groups = defaultdict(list)
    for drug in inj_schedule:
        rp_groups[str(drug.get('投与順序',''))].append(drug)
    sorted_rps = sorted(rp_groups.keys(),
                        key=lambda x: int(x) if x.isdigit() else 99)

    all_invest_days = set()
    for rp in sorted_rps:
        all_invest_days.update(
            parse_days_num(rp_groups[rp][0].get('投与Day数値',''))
        )
    invest_days = sorted(all_invest_days)
    columns = []; prev = 0
    for d in invest_days:
        if d > prev+1:
            rs=prev+1; re=d-1
            columns.append({'type':'休薬中',
                             'label':f'{rs}日目' if rs==re else f'{rs}〜{re}日目',
                             'days':list(range(rs,re+1))})
        columns.append({'type':'投与','label':f'{d}日目','days':[d]})
        prev = d
    if prev < cycle:
        rs=prev+1; re=cycle
        columns.append({'type':'休薬','label':f'〜{re}日目',
                        'days':list(range(rs,re+1))})

    need_bw  = any(str(d.get('用量根拠','')) in ('BW依存','AUC依存')
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')
    need_bsa = any(str(d.get('用量根拠',''))=='BSA依存'
                   for d in drugs if str(d.get('①O欄_抗がん剤',''))=='○')

    rp_all_groups = defaultdict(list)
    for drug in drugs:
        rp = str(drug.get('投与順序',''))
        if rp != '内服':
            rp_all_groups[rp].append(drug)

    total_min_raw = sum(
        max((safe_float(d.get('投与時間数値', 0))
             for d in rp_drugs), default=0)
        for rp_drugs in rp_all_groups.values()
    ) * 60
    # 5分単位で切り上げ
    import math
    total_min = int(math.ceil(total_min_raw / 5) * 5)

    def get_rp_info(rp_drugs):
        rp_sorted = sorted(rp_drugs,
                           key=lambda d: INJECTION_ORDER.get(str(d.get('支持療法分類','')),99))
        names = [str(d.get('商品名','') or d.get('採用商品名（全角）','')) for d in rp_sorted]
        name_text = '\n'.join([n for n in names if n])
        dose_parts = []
        for d in rp_sorted:
            dose_base = str(d.get('用量根拠',''))
            kubun     = str(d.get('薬剤区分',''))
            try:
                _rv3 = str(d.get('投与量数値', 0) or 0).strip()
                _rv3 = ''.join(c for c in _rv3 if c.isdigit() or c == '.')
                dose_num = float(_rv3 or 0)
            except:
                dose_num = 0
            dose_str, unit_str = format_dose_text(d)
            if kubun == '抗がん剤':
                if dose_base=='BSA依存':   dose_parts.append(f'({dose_num}mg/m²)')
                elif dose_base=='AUC依存': dose_parts.append(f'(AUC：{int(dose_num)})')
                elif dose_base=='BW依存':  dose_parts.append(f'({dose_num}mg/kg)')
                elif dose_base=='固定用量':dose_parts.append(f'({dose_str}{unit_str})')
        time_val  = rp_sorted[0].get('投与時間文字','')
        time_text = f'(投与時間：{time_val})' if time_val else ''
        type_str  = ''
        for d in rp_sorted:
            t = str(master_dict.get(str(d.get('管理コード','')),{}).get('スケジュールシール用種類',''))
            if t: type_str = t; break
        drug_days = parse_days_num(rp_sorted[0].get('投与Day数値',''))
        return name_text, dose_parts, time_text, type_str, drug_days

    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mg=Cm(0.8); cw=A4_W-mg*2; x0=mg; y0=mg
    title_w=cw*0.72; date_w=cw*0.28
    title_h=Cm(1.5); course_h=Cm(1.0)

    add_textbox(slide,x0,y0,title_w,title_h,
                text=basic['レジメン名'],
                font_size=Pt(24),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0,y0+title_h,title_w,course_h,
                text=f'（1コース{cycle}日）',
                font_size=Pt(20),bold=True,color=COLOR_BLACK,align=PP_ALIGN.CENTER)
    add_textbox(slide,x0+title_w,y0,date_w,title_h+course_h,
                text='当レジメンの開始日\n年　　月　　日',
                font_size=Pt(11),bold=False,color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    y_cur = y0+title_h+course_h+Cm(0.2)
    note_lines = (['・'+str(n['注意事項文章']) for n in seal_notes] if seal_notes else [
        '・薬剤の他、副作用を予防する、投与ラインを満たす、薬剤を洗い流すために生理食塩液等を投与します。',
        '・強い副作用が出た場合や患者様の状態によっては、治療を延期したり、減量することがあります。'
    ])
    note_h = Cm(0.55*len(note_lines))
    add_textbox_multi(slide,x0,y_cur,cw,note_h,
                      lines=note_lines,font_size=Pt(11),
                      color=COLOR_BLACK,align=PP_ALIGN.LEFT)
    y_cur += note_h+Cm(0.2)

    has_right   = bool(oral_schedule) or need_bw or need_bsa
    tbl_w_ratio = 0.60 if has_right else 1.0
    right_w     = cw*0.38
    right_x     = x0+cw*tbl_w_ratio+cw*0.02

    type_w_cm=2.2; name_w_cm=3.8; inv_w_cm=1.1; rest_w_cm=1.5
    col_w_cm = [inv_w_cm if c['type']=='投与'
                else inv_w_cm*0.8 if c['type']=='休薬中'
                else rest_w_cm for c in columns]
    tbl_total_w  = type_w_cm+name_w_cm+sum(col_w_cm)
    scale        = float((cw*tbl_w_ratio)/Cm(tbl_total_w))
    type_w_cm   *= scale; name_w_cm *= scale
    col_w_cm     = [w*scale for w in col_w_cm]

    n_cols=2+len(columns); n_rps=len(sorted_rps); n_rows=1+n_rps

    def calc_lines(nd,dp,tt):
        return nd.count('\n')+1+len(dp)+(1 if tt else 0)

    row_heights = [Cm(0.65)]
    for rp in sorted_rps:
        nd,dp,tt,_,_ = get_rp_info(rp_groups[rp])
        row_heights.append(max(Cm(0.75),Cm(0.42*calc_lines(nd,dp,tt)+0.2)))
    tbl_h = sum(row_heights)

    shape = slide.shapes.add_table(n_rows,n_cols,x0,y_cur,
                                   Cm(tbl_total_w*scale),tbl_h)
    table = shape.table
    table.columns[0].width = Cm(type_w_cm)
    table.columns[1].width = Cm(name_w_cm)
    for i,w in enumerate(col_w_cm): table.columns[2+i].width = Cm(w)
    for i,h in enumerate(row_heights): table.rows[i].height = h

    for col_i,label in enumerate(['種類','薬剤名']):
        cell = table.cell(0,col_i)
        set_cell_text_multi(cell,[label],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)
    for i,col in enumerate(columns):
        cell = table.cell(0,2+i)
        set_cell_text_multi(cell,[col['label']],font_size=Pt(12),bold=True,
                             color=COLOR_WHITE,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,COLOR_HEADER); set_cell_margin(cell)

    rest_col_indices = [(i,col) for i,col in enumerate(columns)
                        if col['type'] in ('休薬','休薬中')]
    for row_i,rp in enumerate(sorted_rps):
        r=row_i+1
        name_text,dose_parts,time_text,type_str,drug_days = get_rp_info(rp_groups[rp])
        row_color = TYPE_COLOR_MAP.get(type_str,COLOR_NAUSEA)
        cell = table.cell(r,0)
        set_cell_text_multi(cell,[type_str],font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        name_lines = name_text.split('\n')+dose_parts
        if time_text: name_lines.append(time_text)
        cell = table.cell(r,1)
        set_cell_text_multi(cell,name_lines,font_size=Pt(11),bold=True,
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_bg(cell,row_color); set_cell_margin(cell)
        for i,col in enumerate(columns):
            cell = table.cell(r,2+i); set_cell_margin(cell)
            if col['type']=='投与':
                hit = any(d in drug_days for d in col['days'])
                if hit:
                    set_cell_text_multi(cell,['●'],font_size=Pt(14),bold=True,
                                         color=COLOR_BLACK,align=PP_ALIGN.CENTER)
                    set_cell_bg(cell,row_color)
                else:
                    set_cell_bg(cell,COLOR_REST_BG)

    for i,col in rest_col_indices:
        col_idx = 2+i
        merge_cells_vertical(table,col_idx,1,n_rps)
        cell = table.cell(1,col_idx)
        set_cell_bg(cell,COLOR_REST_BG)
        set_cell_text_multi(cell,[''],font_size=Pt(11),
                             color=COLOR_BLACK,align=PP_ALIGN.CENTER)
        set_cell_margin(cell)

    tbl_bottom  = y_cur+tbl_h
    y_after_tbl = tbl_bottom+Cm(0.15)
    if tegaki_notes:
        nb_h = Cm(0.5*len(tegaki_notes))
        add_textbox_multi(slide,x0,y_after_tbl,cw*tbl_w_ratio,nb_h,
                          lines=[str(n['注意事項文章']) for n in tegaki_notes],
                          font_size=Pt(11),color=COLOR_BLACK,align=PP_ALIGN.LEFT)

    if has_right:
        ry = y_cur
        if oral_schedule:
            oral_lines = []
            for drug in oral_schedule:
                name   = str(drug.get('商品名','') or drug.get('採用商品名（全角）',''))
                timing = str(drug.get('投与タイミング',''))
                line   = f'・{name}'
                if timing and timing not in ('','ー','nan'):
                    line += f'\n　{timing}'
                oral_lines.append(line)
            oral_box_h = Cm(0.55*sum(l.count('\n')+1 for l in oral_lines)+0.3)
            add_textbox_multi(slide,right_x+Cm(0.15),ry+Cm(0.1),
                              right_w-Cm(0.3),oral_box_h,
                              lines=oral_lines,font_size=Pt(11),
                              color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += oral_box_h+Cm(0.3)
        if need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(0.7),
                        text='体重　　　　　kg',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(0.9)
        if need_bsa or need_bw:
            add_textbox(slide,right_x,ry,right_w,Cm(1.0),
                        text='体表面積\n\n＿＿＿＿m²',font_size=Pt(11),
                        color=COLOR_BLACK,align=PP_ALIGN.LEFT)
            ry += Cm(1.3)
        add_textbox(slide,right_x,ry,right_w,Cm(0.8),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        ry += Cm(1.0)
        # 短縮注記を表示（対象薬剤が含まれる場合）
        shorten_notes = []
        for _d in drugs:
            _sn = str(_d.get('短縮注記','')).strip()
            if _sn and _sn not in shorten_notes:
                shorten_notes.append(_sn)
        if shorten_notes:
            for sn in shorten_notes:
                add_textbox(slide,right_x,ry,right_w,Cm(0.6),
                            text=sn,
                            font_size=Pt(9),bold=False,
                            color=RGBColor(0x85,0x64,0x04),
                            align=PP_ALIGN.LEFT)
                ry += Cm(0.65)
        add_textbox(slide,right_x,ry,right_w,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
    else:
        footer_y = A4_H-Cm(1.5)-mg
        add_textbox(slide,x0,footer_y,cw*0.5,Cm(0.7),
                    text='体表面積　　　　m²',font_size=Pt(11),
                    color=COLOR_BLACK,align=PP_ALIGN.LEFT)
        add_textbox(slide,x0+cw*0.5,footer_y,cw*0.5,Cm(0.7),
                    text=f'所要時間の目安：{total_min}分',
                    font_size=Pt(14),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)
        # 短縮注記を表示（対象薬剤が含まれる場合）
        shorten_notes_b = []
        for _d in drugs:
            _sn = str(_d.get('短縮注記','')).strip()
            if _sn and _sn not in shorten_notes_b:
                shorten_notes_b.append(_sn)
        if shorten_notes_b:
            for sn in shorten_notes_b:
                add_textbox(slide,x0,footer_y+Cm(0.8),cw,Cm(0.5),
                            text=sn,
                            font_size=Pt(9),bold=False,
                            color=RGBColor(0x85,0x64,0x04),
                            align=PP_ALIGN.LEFT)
                footer_y += Cm(0.55)
        add_textbox(slide,x0,footer_y+Cm(0.8),cw,Cm(0.6),
                    text='東北大学病院　薬剤部',
                    font_size=Pt(12),bold=True,
                    color=COLOR_BLACK,align=PP_ALIGN.RIGHT)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()



